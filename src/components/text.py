from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from src.components._style import style_runs, set_paragraph_text


def add_title(slide, theme, text, left, top, width=None, subtitle=None):
    """スライドタイトル: 左上配置、太字、下に区切り線。
    subtitle を渡すと区切り線直下にサブヘッド (小さいグレー 1 行) を追加。"""
    if width is None:
        width = theme.content_width

    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    set_paragraph_text(
        tf.paragraphs[0], text,
        size=theme.font_size_h1, color=theme.text_primary,
        name=theme.font_title, bold=True,
    )

    line_top = top + Inches(0.6)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, line_top, width, Pt(2)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = theme.secondary
    line.line.fill.background()

    if subtitle:
        sub_box = slide.shapes.add_textbox(
            left, line_top + Inches(0.05), width, Inches(0.35),
        )
        sub_tf = sub_box.text_frame
        sub_tf.word_wrap = True
        set_paragraph_text(
            sub_tf.paragraphs[0], subtitle,
            size=theme.font_size_h3, color=theme.text_secondary,
            name=theme.font_body,
        )


def add_subtitle(slide, theme, text, left, top, width=None):
    """サブタイトル: タイトル直下、やや小さく"""
    if width is None:
        width = theme.content_width

    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    set_paragraph_text(
        tf.paragraphs[0], text,
        size=theme.font_size_h2, color=theme.text_secondary, name=theme.font_body,
    )


def add_bullets(slide, theme, items, left, top, width=None, height=None):
    """箇条書き: インデント2階層対応、行頭は控えめ記号"""
    if width is None:
        width = theme.content_width
    if height is None:
        height = Inches(4.0)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(item, dict):
            text = item["text"]
            level = item.get("level", 0)
        else:
            text = item
            level = 0

        prefix = "  " * level + "― "
        p.text = prefix + text
        color = theme.text_secondary if level > 0 else theme.text_primary
        style_runs(p, size=theme.font_size_body, color=color, name=theme.font_body)
        p.space_after = Pt(6)
        if level > 0:
            p.level = level


def add_callout(slide, theme, text, left, top, width=None, height=None, variant="info"):
    """強調ボックス: 左アクセントバー付き矩形内にテキスト。
    variant: "info" (デフォルト/primary) | "success" | "warning" | "danger" """
    if width is None:
        width = theme.content_width
    if height is None:
        line_count = text.count("\n") + 1
        height = Inches(0.3 + 0.3 * line_count)

    color_map = {
        "info": theme.primary,
        "success": theme.success,
        "warning": theme.warning,
        "danger": theme.danger,
    }
    accent = color_map.get(variant, theme.primary)

    r, g, b = accent[0], accent[1], accent[2]
    light_r = r + (255 - r) * 9 // 10
    light_g = g + (255 - g) * 9 // 10
    light_b = b + (255 - b) * 9 // 10

    # Light background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(light_r, light_g, light_b)
    bg.line.fill.background()

    # Left accent bar
    bar_w = Inches(0.07)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, bar_w, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    tf = bg.text_frame
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    tf.word_wrap = True
    set_paragraph_text(
        tf.paragraphs[0], text,
        size=theme.font_size_body, color=theme.text_primary,
        name=theme.font_body, bold=True,
    )


def add_footnote(slide, theme, text, left, bottom_margin=None):
    """脚注: スライド下部、小フォント"""
    if bottom_margin is None:
        bottom_margin = theme.margin_bottom

    top = theme.slide_height - bottom_margin - Inches(0.3)
    width = theme.content_width

    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    set_paragraph_text(
        tf.paragraphs[0], text,
        size=theme.font_size_caption, color=theme.text_secondary, name=theme.font_body,
    )
