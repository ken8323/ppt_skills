from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN


def add_timeline(slide, theme, milestones, left, top, width=None, height=None, today=None):
    """タイムライン: 横軸に時間、マイルストーンを上下にプロット
    milestones: [{"date": "2026/4", "label": "キックオフ"}, ...]
    today: str — 任意。日付文字列 (例 "2026Q3")。該当位置に「現在」縦マーカーを表示。
    """
    if width is None:
        width = theme.content_width
    if height is None:
        height = Inches(3.0)

    n = len(milestones)

    line_y = top + height // 2
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, line_y, width, Pt(3),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = theme.primary
    line.line.fill.background()

    for i, ms in enumerate(milestones):
        x = left + (width * i) // (n - 1) if n > 1 else left + width // 2
        dot_size = Inches(0.25)

        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x - dot_size // 2, line_y - dot_size // 2 + Pt(1),
            dot_size, dot_size,
        )
        color_idx = i % len(theme.chart_colors)
        dot.fill.solid()
        dot.fill.fore_color.rgb = theme.chart_colors[color_idx]
        dot.line.fill.background()

        is_above = (i % 2 == 0)
        label_width = Inches(1.5)
        label_x = x - label_width // 2

        if is_above:
            date_top = line_y - Inches(1.2)
            label_top = line_y - Inches(0.8)
        else:
            date_top = line_y + Inches(0.5)
            label_top = line_y + Inches(0.9)

        date_box = slide.shapes.add_textbox(label_x, date_top, label_width, Inches(0.3))
        date_box.text_frame.paragraphs[0].text = ms["date"]
        date_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        for run in date_box.text_frame.paragraphs[0].runs:
            run.font.size = theme.font_size_caption
            run.font.color.rgb = theme.text_secondary
            run.font.name = theme.font_body
            run.font.bold = True

        label_box = slide.shapes.add_textbox(label_x, label_top, label_width, Inches(0.4))
        label_box.text_frame.word_wrap = True
        label_box.text_frame.paragraphs[0].text = ms["label"]
        label_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        for run in label_box.text_frame.paragraphs[0].runs:
            run.font.size = theme.font_size_body
            run.font.color.rgb = theme.text_primary
            run.font.name = theme.font_body

    if today and n > 1:
        _add_today_marker(slide, theme, today, milestones, left, top, width, height, line_y)


def _add_today_marker(slide, theme, today, milestones, left, top, width, height, line_y):
    """現在位置マーカー: 日付文字列の辞書順でマイルストーン間を補間"""
    dates = [ms["date"] for ms in milestones]
    n = len(dates)

    # 補間位置を 0.0 ~ (n-1) の浮動小数で求める
    if today <= dates[0]:
        pos = 0.0
    elif today >= dates[-1]:
        pos = float(n - 1)
    else:
        for i in range(1, n):
            if today < dates[i]:
                pos = i - 1 + 0.5  # 前後のちょうど中間
                break
        else:
            pos = float(n - 1)

    today_x = left + int(width * pos / (n - 1))
    marker_h = height

    marker = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        today_x - Pt(1), top, Pt(2), marker_h,
    )
    marker.fill.solid()
    marker.fill.fore_color.rgb = theme.danger
    marker.line.fill.background()

    label_box = slide.shapes.add_textbox(
        today_x - Inches(0.4), top - Inches(0.28),
        Inches(0.8), Inches(0.25),
    )
    p = label_box.text_frame.paragraphs[0]
    p.text = "現在"
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        run.font.size = theme.font_size_caption
        run.font.color.rgb = theme.danger
        run.font.bold = True
        run.font.name = theme.font_body


def add_gantt(slide, theme, tasks, phases, left, top, width=None, height=None):
    """ガントチャート: 横棒で期間表示、フェーズ色分け
    tasks: [{"name": "要件定義", "start": 0, "duration": 2, "progress": 0.6}, ...]
    phases: ["4月", "5月", "6月", ...] ヘッダラベル
    progress: 0.0-1.0 — 任意。指定するとバーの完了比率を濃色で塗り潰す。
    """
    if width is None:
        width = theme.content_width
    if height is None:
        height = Inches(4.0)

    n_tasks = len(tasks)
    n_phases = len(phases)

    label_width = Inches(2.0)
    chart_width = width - label_width
    phase_width = chart_width // n_phases
    row_height = min(Inches(0.6), height // (n_tasks + 1))

    for j, phase_label in enumerate(phases):
        hdr_left = left + label_width + phase_width * j
        hdr = slide.shapes.add_textbox(hdr_left, top, phase_width, row_height)
        hdr.text_frame.paragraphs[0].text = phase_label
        hdr.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        for run in hdr.text_frame.paragraphs[0].runs:
            run.font.size = theme.font_size_caption
            run.font.color.rgb = theme.text_secondary
            run.font.name = theme.font_body
            run.font.bold = True

    hdr_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left + label_width, top + row_height - Pt(1),
        chart_width, Pt(1),
    )
    hdr_line.fill.solid()
    hdr_line.fill.fore_color.rgb = theme.border
    hdr_line.line.fill.background()

    for i, task in enumerate(tasks):
        row_top = top + row_height * (i + 1) + Inches(0.1)

        name_box = slide.shapes.add_textbox(left, row_top, label_width, row_height)
        name_box.text_frame.paragraphs[0].text = task["name"]
        for run in name_box.text_frame.paragraphs[0].runs:
            run.font.size = theme.font_size_body
            run.font.color.rgb = theme.text_primary
            run.font.name = theme.font_body

        bar_left = left + label_width + phase_width * task["start"]
        bar_width = phase_width * task["duration"]
        bar_height = row_height - Inches(0.15)

        bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            bar_left, row_top + Inches(0.05), bar_width, bar_height,
        )
        color_idx = i % len(theme.chart_colors)
        bar.fill.solid()
        bar.fill.fore_color.rgb = theme.chart_colors[color_idx]
        bar.line.fill.background()

        tf = bar.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].text = task["name"]
        for run in tf.paragraphs[0].runs:
            run.font.size = theme.font_size_caption
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            run.font.name = theme.font_body
            run.font.bold = True

        progress = task.get("progress", 0.0)
        if progress and progress > 0:
            done_width = max(int(bar_width * progress), Inches(0.1))
            done_bar = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                bar_left, row_top + Inches(0.05), done_width, bar_height,
            )
            c = theme.chart_colors[color_idx]
            dark_r = c[0] * 65 // 100
            dark_g = c[1] * 65 // 100
            dark_b = c[2] * 65 // 100
            done_bar.fill.solid()
            done_bar.fill.fore_color.rgb = RGBColor(dark_r, dark_g, dark_b)
            done_bar.line.fill.background()
