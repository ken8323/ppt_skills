from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN


ICON_SHAPES = {
    "circle": MSO_SHAPE.OVAL,
    "square": MSO_SHAPE.ROUNDED_RECTANGLE,
    "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
    "diamond": MSO_SHAPE.DIAMOND,
    "star": MSO_SHAPE.STAR_5_POINT,
    "check": MSO_SHAPE.OVAL,
    "arrow_right": MSO_SHAPE.RIGHT_ARROW,
    "arrow_up": MSO_SHAPE.UP_ARROW,
    "hexagon": MSO_SHAPE.HEXAGON,
    "lightning": MSO_SHAPE.LIGHTNING_BOLT,
}

ICON_SYMBOLS = {
    "check": "✓",
    "arrow_right": "→",
    "arrow_up": "↑",
    "circle": "",
    "square": "",
    "triangle": "",
    "diamond": "",
    "star": "",
    "hexagon": "",
    "lightning": "",
}


def add_icon_with_label(slide, theme, icon_type, label, left, top, size=None, color_idx=0):
    """アイコン+ラベル: 丸/角丸内に幾何学記号+下にテキスト"""
    if size is None:
        size = Inches(0.8)

    shape_type = ICON_SHAPES.get(icon_type, MSO_SHAPE.OVAL)
    symbol = ICON_SYMBOLS.get(icon_type, "")

    icon = slide.shapes.add_shape(shape_type, left, top, size, size)
    c_idx = color_idx % len(theme.chart_colors)
    icon.fill.solid()
    icon.fill.fore_color.rgb = theme.chart_colors[c_idx]
    icon.line.fill.background()

    if symbol:
        tf = icon.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].text = symbol
        for run in tf.paragraphs[0].runs:
            run.font.size = Pt(int(size / Inches(1) * 18))
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            run.font.bold = True

    label_box = slide.shapes.add_textbox(
        left - Inches(0.3), top + size + Inches(0.1),
        size + Inches(0.6), Inches(0.4),
    )
    label_box.text_frame.word_wrap = True
    label_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    label_box.text_frame.paragraphs[0].text = label
    for run in label_box.text_frame.paragraphs[0].runs:
        run.font.size = theme.font_size_caption
        run.font.color.rgb = theme.text_primary
        run.font.name = theme.font_body


def add_kpi_card(slide, theme, value, unit, label, left, top, width=None, height=None, color_idx=0):
    """KPI表示: 大きな数字+単位+ラベル、カード型"""
    if width is None:
        width = Inches(3.0)
    if height is None:
        height = Inches(2.0)

    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
    )
    card.fill.solid()
    c = theme.chart_colors[color_idx % len(theme.chart_colors)]
    light_r = c[0] + (255 - c[0]) * 92 // 100
    light_g = c[1] + (255 - c[1]) * 92 // 100
    light_b = c[2] + (255 - c[2]) * 92 // 100
    card.fill.fore_color.rgb = RGBColor(light_r, light_g, light_b)
    card.line.color.rgb = theme.border
    card.line.width = Pt(0.5)

    tf = card.text_frame
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.2)
    tf.word_wrap = True

    p_value = tf.paragraphs[0]
    p_value.alignment = PP_ALIGN.CENTER
    run_val = p_value.add_run()
    run_val.text = value
    run_val.font.size = Pt(36)
    run_val.font.bold = True
    run_val.font.color.rgb = theme.text_primary
    run_val.font.name = theme.font_title

    run_unit = p_value.add_run()
    run_unit.text = " " + unit
    run_unit.font.size = Pt(16)
    run_unit.font.color.rgb = theme.text_secondary
    run_unit.font.name = theme.font_body

    p_label = tf.add_paragraph()
    p_label.alignment = PP_ALIGN.CENTER
    p_label.space_before = Pt(8)
    run_label = p_label.add_run()
    run_label.text = label
    run_label.font.size = theme.font_size_body
    run_label.font.color.rgb = theme.text_secondary
    run_label.font.name = theme.font_body


def add_icon_row(slide, theme, items, left, top, width=None, icon_size=None):
    """アイコン横並び: 3-5個のアイコン+ラベルを等間隔配置
    items: [{"icon": "circle", "label": "項目1"}, ...]
    """
    if width is None:
        width = theme.content_width
    if icon_size is None:
        icon_size = Inches(0.8)

    n = len(items)
    spacing = width // n

    for i, item in enumerate(items):
        icon_left = left + spacing * i + (spacing - icon_size) // 2
        add_icon_with_label(
            slide, theme,
            item["icon"], item["label"],
            icon_left, top,
            size=icon_size,
            color_idx=i,
        )
