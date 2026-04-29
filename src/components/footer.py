"""スライド下端のフッター + ページ番号。generator から自動注入される。
cover / section_divider / thank_you 系の「背景塗りつぶしページ」ではスキップする
（レイアウトごとに skip_footer フラグで制御）。"""

from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

from src.components._style import set_paragraph_text


FOOTER_HEIGHT = Inches(0.28)


def add_page_footer(slide, theme, page_num, total, footer_text=""):
    """左下に brand_name/footer_text、右下に「n / N」。上端に細いディバイダライン。"""
    sw = theme.slide_width
    sh = theme.slide_height

    divider_top = sh - theme.margin_bottom + Inches(0.05)
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        theme.margin_left, divider_top,
        sw - theme.margin_left - theme.margin_right, Pt(0.5),
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = theme.border
    divider.line.fill.background()

    footer_top = divider_top + Inches(0.08)

    left_text = footer_text or theme.brand_name
    if left_text:
        left_box = slide.shapes.add_textbox(
            theme.margin_left, footer_top,
            Inches(8.0), FOOTER_HEIGHT,
        )
        left_tf = left_box.text_frame
        left_tf.margin_top = 0
        left_tf.margin_bottom = 0
        set_paragraph_text(
            left_tf.paragraphs[0], left_text,
            size=theme.font_size_footnote, color=theme.text_secondary, name=theme.font_body,
        )

    right_box = slide.shapes.add_textbox(
        sw - theme.margin_right - Inches(2.0), footer_top,
        Inches(2.0), FOOTER_HEIGHT,
    )
    right_tf = right_box.text_frame
    right_tf.margin_top = 0
    right_tf.margin_bottom = 0
    right_tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    set_paragraph_text(
        right_tf.paragraphs[0], f"{page_num} / {total}",
        size=theme.font_size_footnote, color=theme.text_secondary, name=theme.font_body,
    )
