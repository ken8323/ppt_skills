import math
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN


def add_matrix_2x2(slide, theme, x_axis, y_axis, quadrants, left, top, size=None,
                   recommended_quadrant=None):
    """2x2マトリクス: 軸ラベル付き、各象限にテキスト配置
    quadrants: [top-left, top-right, bottom-left, bottom-right]
    recommended_quadrant: 0-3 で指定すると、その象限を太枠で強調
    """
    if size is None:
        size = Inches(5.0)

    cell_size = size // 2
    gap = Inches(0.05)

    colors = [
        theme.chart_colors[0],
        theme.chart_colors[2],
        theme.chart_colors[3],
        theme.chart_colors[1],
    ]

    positions = [
        (left, top),
        (left + cell_size + gap, top),
        (left, top + cell_size + gap),
        (left + cell_size + gap, top + cell_size + gap),
    ]

    for i, (qtext, (ql, qt)) in enumerate(zip(quadrants, positions)):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, ql, qt, cell_size - gap, cell_size - gap
        )
        shape.fill.solid()
        c = colors[i]
        light_r = c[0] + (255 - c[0]) * 8 // 10
        light_g = c[1] + (255 - c[1]) * 8 // 10
        light_b = c[2] + (255 - c[2]) * 8 // 10
        shape.fill.fore_color.rgb = RGBColor(light_r, light_g, light_b)

        if i == recommended_quadrant:
            shape.line.color.rgb = theme.primary
            shape.line.width = Pt(2.5)
        else:
            shape.line.color.rgb = theme.border
            shape.line.width = Pt(0.5)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].text = qtext
        for run in tf.paragraphs[0].runs:
            run.font.size = theme.font_size_body
            run.font.name = theme.font_body
            run.font.color.rgb = theme.text_primary

    x_label = slide.shapes.add_textbox(
        left, top + size + Inches(0.1), size, Inches(0.3)
    )
    x_label.text_frame.paragraphs[0].text = x_axis
    x_label.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    for run in x_label.text_frame.paragraphs[0].runs:
        run.font.size = theme.font_size_caption
        run.font.color.rgb = theme.text_secondary
        run.font.name = theme.font_body

    y_label = slide.shapes.add_textbox(
        left - Inches(0.8), top + size // 2 - Inches(0.15), Inches(0.7), Inches(0.3)
    )
    y_label.text_frame.paragraphs[0].text = y_axis
    y_label.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    for run in y_label.text_frame.paragraphs[0].runs:
        run.font.size = theme.font_size_caption
        run.font.color.rgb = theme.text_secondary
        run.font.name = theme.font_body


def add_pyramid(slide, theme, levels, left, top, width=None, height=None):
    """ピラミッド: 台形の積み重ね
    levels: list[str] または list[{"text": str, "note": str}]
    note を指定すると右側に1行注釈を表示する。
    """
    if width is None:
        width = Inches(6.0)
    if height is None:
        height = Inches(4.5)

    n = len(levels)
    level_height = height // n
    gap = Inches(0.03)

    for i, level in enumerate(levels):
        if isinstance(level, dict):
            text = level.get("text", "")
            note = level.get("note")
        else:
            text = level
            note = None

        level_width = int(width * (0.3 + 0.7 * (i + 1) / n))
        level_left = left + (width - level_width) // 2

        shape = slide.shapes.add_shape(
            MSO_SHAPE.TRAPEZOID,
            level_left, top + level_height * i + gap * i,
            level_width, level_height - gap,
        )
        color_idx = i % len(theme.chart_colors)
        shape.fill.solid()
        shape.fill.fore_color.rgb = theme.chart_colors[color_idx]
        shape.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        shape.line.width = Pt(2)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].text = text
        for run in tf.paragraphs[0].runs:
            run.font.size = theme.font_size_body
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            run.font.name = theme.font_title
            run.font.bold = True

        if note:
            note_top = top + level_height * i + (level_height - Inches(0.3)) // 2
            note_box = slide.shapes.add_textbox(
                left + width + Inches(0.2), note_top,
                Inches(2.5), Inches(0.3),
            )
            note_box.text_frame.paragraphs[0].text = note
            for run in note_box.text_frame.paragraphs[0].runs:
                run.font.size = theme.font_size_caption
                run.font.color.rgb = theme.text_secondary
                run.font.name = theme.font_body


def add_process_flow(slide, theme, steps, left, top, width=None, height=None, style="arrow"):
    """プロセスフロー: 矢印で繋がった角丸矩形、横並び
    style: "arrow" (デフォルト) | "chevron" — chevron は山形のステップ連結
    """
    if width is None:
        width = theme.content_width
    if height is None:
        height = Inches(1.5)

    n = len(steps)

    if style == "chevron":
        overlap = Inches(0.15)
        cell_w = width // n
        for i, step_text in enumerate(steps):
            shape_x = left + cell_w * i
            shape_w = cell_w + (overlap if i < n - 1 else 0)
            shape_type = MSO_SHAPE.PENTAGON if i == 0 else MSO_SHAPE.CHEVRON
            shape = slide.shapes.add_shape(shape_type, shape_x, top, shape_w, height)
            color_idx = i % len(theme.chart_colors)
            shape.fill.solid()
            shape.fill.fore_color.rgb = theme.chart_colors[color_idx]
            shape.line.fill.background()
            tf = shape.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.2 if i > 0 else 0.08)
            tf.margin_right = Inches(0.1)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.text = step_text
            for run in p.runs:
                run.font.size = Pt(13)
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                run.font.name = theme.font_title
                run.font.bold = True
        return

    # --- arrow style (デフォルト) ---
    arrow_width = Inches(0.25)
    total_arrow_width = arrow_width * (n - 1)

    # ボックス幅を上限2.2"に抑えて中央寄せ
    max_box_width = Inches(2.2)
    box_width = min((width - total_arrow_width) // n, max_box_width)
    effective_width = box_width * n + total_arrow_width
    flow_left = left + (width - effective_width) // 2
    box_height = height

    for i, step_text in enumerate(steps):
        box_left = flow_left + i * (box_width + arrow_width)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            box_left, top, box_width, box_height,
        )
        color_idx = i % len(theme.chart_colors)
        shape.fill.solid()
        shape.fill.fore_color.rgb = theme.chart_colors[color_idx]
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = step_text
        for run in p.runs:
            run.font.size = Pt(13)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            run.font.name = theme.font_title
            run.font.bold = True

        if i < n - 1:
            arrow_left = box_left + box_width
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                arrow_left, top + box_height // 2 - Inches(0.15),
                arrow_width, Inches(0.3),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = theme.text_secondary
            arrow.line.fill.background()


def add_cycle(slide, theme, items, left, top, size=None):
    """サイクル図: 円形配置の要素"""
    if size is None:
        size = Inches(5.0)

    n = len(items)
    center_x = left + size // 2
    center_y = top + size // 2
    radius = size // 2 - Inches(0.6)
    node_size = Inches(1.4)

    for i, text in enumerate(items):
        angle = -math.pi / 2 + (2 * math.pi * i / n)
        node_x = int(center_x + radius * math.cos(angle) - node_size // 2)
        node_y = int(center_y + radius * math.sin(angle) - node_size // 2)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, node_x, node_y, node_size, node_size,
        )
        color_idx = i % len(theme.chart_colors)
        shape.fill.solid()
        shape.fill.fore_color.rgb = theme.chart_colors[color_idx]
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].text = text
        for run in tf.paragraphs[0].runs:
            run.font.size = Pt(12)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            run.font.name = theme.font_title
            run.font.bold = True

        mid_angle = angle + math.pi / n
        arrow_x = int(center_x + (radius + Inches(0.3)) * math.cos(mid_angle) - Inches(0.15))
        arrow_y = int(center_y + (radius + Inches(0.3)) * math.sin(mid_angle) - Inches(0.15))

        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, arrow_x, arrow_y, Inches(0.3), Inches(0.3),
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = theme.text_secondary
        arrow.line.fill.background()
        arrow.rotation = math.degrees(mid_angle + math.pi / 2)


def add_org_chart(slide, theme, data, left, top, width=None, height=None):
    """組織図: ツリー構造、線で接続
    data: {"name": "CEO", "children": [{"name": "CTO"}, ...]}
    """
    if width is None:
        width = theme.content_width
    if height is None:
        height = Inches(4.5)

    levels = []
    _collect_levels(data, 0, levels)

    level_height = height // max(len(levels), 1)
    node_height = Inches(0.6)

    _render_org_node(slide, theme, data, left, top, width, level_height, node_height)


def _collect_levels(node, depth, levels):
    while len(levels) <= depth:
        levels.append(0)
    levels[depth] += 1
    for child in node.get("children", []):
        _collect_levels(child, depth + 1, levels)


def _render_org_node(slide, theme, node, area_left, area_top, area_width,
                     level_height, node_height, depth=0):
    node_width = Inches(2.0)
    node_left = area_left + (area_width - node_width) // 2
    node_top = area_top

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, node_left, node_top, node_width, node_height,
    )
    if depth == 0:
        shape.fill.solid()
        shape.fill.fore_color.rgb = theme.primary
        text_color = RGBColor(0xFF, 0xFF, 0xFF)
    else:
        shape.fill.solid()
        color_idx = depth % len(theme.chart_colors)
        c = theme.chart_colors[color_idx]
        light_r = c[0] + (255 - c[0]) * 7 // 10
        light_g = c[1] + (255 - c[1]) * 7 // 10
        light_b = c[2] + (255 - c[2]) * 7 // 10
        shape.fill.fore_color.rgb = RGBColor(light_r, light_g, light_b)
        text_color = theme.text_primary
    shape.line.color.rgb = theme.border
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].text = node["name"]
    for run in tf.paragraphs[0].runs:
        run.font.size = Pt(12)
        run.font.color.rgb = text_color
        run.font.name = theme.font_body
        if depth == 0:
            run.font.bold = True

    children = node.get("children", [])
    if children:
        n = len(children)
        child_area_width = area_width // n

        for i, child in enumerate(children):
            child_left = area_left + child_area_width * i
            child_top = area_top + level_height

            parent_center_x = node_left + node_width // 2
            child_center_x = child_left + child_area_width // 2

            mid_y = area_top + node_height + (level_height - node_height) // 2
            line1 = slide.shapes.add_connector(
                1, parent_center_x, node_top + node_height,
                parent_center_x, mid_y,
            )
            line1.line.color.rgb = theme.border
            line1.line.width = Pt(1.5)

            line2 = slide.shapes.add_connector(
                1, parent_center_x, mid_y, child_center_x, mid_y,
            )
            line2.line.color.rgb = theme.border
            line2.line.width = Pt(1.5)

            line3 = slide.shapes.add_connector(
                1, child_center_x, mid_y, child_center_x, child_top,
            )
            line3.line.color.rgb = theme.border
            line3.line.width = Pt(1.5)

            _render_org_node(
                slide, theme, child, child_left, child_top,
                child_area_width, level_height, node_height, depth + 1,
            )


def add_pillars(slide, theme, items, left, top, width=None, height=None):
    """3-5本柱: 各カラムに色付き見出し+本文+任意KPI数値
    items: [{"title": str, "body": str, "kpi": str?}]
    """
    if width is None:
        width = theme.content_width
    if height is None:
        height = Inches(4.0)

    n = len(items)
    gap = Inches(0.2)
    col_width = (width - gap * (n - 1)) // n
    header_height = Inches(0.55)

    for i, item in enumerate(items):
        col_left = left + i * (col_width + gap)
        color_idx = i % len(theme.chart_colors)
        col_color = theme.chart_colors[color_idx]

        header = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, col_left, top, col_width, header_height,
        )
        header.fill.solid()
        header.fill.fore_color.rgb = col_color
        header.line.fill.background()
        tf = header.text_frame
        tf.margin_top = Inches(0.05)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = item.get("title", "")
        for run in p.runs:
            run.font.size = theme.font_size_h3
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            run.font.name = theme.font_title

        body_top = top + header_height + Inches(0.05)
        body_height = height - header_height - Inches(0.05)
        c = col_color
        light_r = c[0] + (255 - c[0]) * 90 // 100
        light_g = c[1] + (255 - c[1]) * 90 // 100
        light_b = c[2] + (255 - c[2]) * 90 // 100

        body_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, col_left, body_top, col_width, body_height,
        )
        body_shape.fill.solid()
        body_shape.fill.fore_color.rgb = RGBColor(light_r, light_g, light_b)
        body_shape.line.color.rgb = col_color
        body_shape.line.width = Pt(0.75)
        tf2 = body_shape.text_frame
        tf2.word_wrap = True
        tf2.margin_left = Inches(0.12)
        tf2.margin_right = Inches(0.12)
        tf2.margin_top = Inches(0.15)

        body_text = item.get("body", "")
        kpi = item.get("kpi")
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        p2.text = body_text
        for run in p2.runs:
            run.font.size = theme.font_size_body
            run.font.color.rgb = theme.text_primary
            run.font.name = theme.font_body

        if kpi:
            p3 = tf2.add_paragraph()
            p3.alignment = PP_ALIGN.CENTER
            p3.space_before = Pt(16)
            r = p3.add_run()
            r.text = kpi
            r.font.size = Pt(30)
            r.font.bold = True
            r.font.color.rgb = col_color
            r.font.name = theme.font_title


def add_swot(slide, theme, cells, left, top, width=None, height=None):
    """SWOT/3C/4P等の2x2フレームワーク枠。軸ラベルなし。
    cells: [{"title": str, "items": list[str]}] — 左上/右上/左下/右下の順
    """
    if width is None:
        width = theme.content_width
    if height is None:
        height = Inches(4.5)

    gap = Inches(0.1)
    cell_w = (width - gap) // 2
    cell_h = (height - gap) // 2

    positions = [
        (left, top),
        (left + cell_w + gap, top),
        (left, top + cell_h + gap),
        (left + cell_w + gap, top + cell_h + gap),
    ]

    for i, (cell, (cl, ct)) in enumerate(zip(cells, positions)):
        col_color = theme.chart_colors[i % len(theme.chart_colors)]
        title_h = Inches(0.4)

        title_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, cl, ct, cell_w, title_h,
        )
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = col_color
        title_shape.line.fill.background()
        tf = title_shape.text_frame
        tf.margin_top = Inches(0.04)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = cell.get("title", "")
        for run in p.runs:
            run.font.size = theme.font_size_h3
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            run.font.name = theme.font_title

        body_h = cell_h - title_h - Inches(0.03)
        c = col_color
        light_r = c[0] + (255 - c[0]) * 88 // 100
        light_g = c[1] + (255 - c[1]) * 88 // 100
        light_b = c[2] + (255 - c[2]) * 88 // 100

        body_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, cl, ct + title_h + Inches(0.03), cell_w, body_h,
        )
        body_shape.fill.solid()
        body_shape.fill.fore_color.rgb = RGBColor(light_r, light_g, light_b)
        body_shape.line.color.rgb = col_color
        body_shape.line.width = Pt(0.5)
        tf2 = body_shape.text_frame
        tf2.word_wrap = True
        tf2.margin_left = Inches(0.12)
        tf2.margin_right = Inches(0.08)
        tf2.margin_top = Inches(0.1)

        for j, item_text in enumerate(cell.get("items", [])):
            p2 = tf2.paragraphs[0] if j == 0 else tf2.add_paragraph()
            p2.text = f"\u2022 {item_text}"
            for run in p2.runs:
                run.font.size = theme.font_size_body
                run.font.color.rgb = theme.text_primary
                run.font.name = theme.font_body


def add_heatmap(slide, theme, col_headers, row_headers, values, left, top,
                width=None, height=None):
    """ヒートマップ: N×M マトリクスに値を色濃淡で表現
    col_headers: list[str], row_headers: list[str]
    values: list[list[float]] — row×col の2次元リスト
    """
    if width is None:
        width = theme.content_width
    if height is None:
        height = Inches(3.5)

    n_rows = len(row_headers)
    n_cols = len(col_headers)
    label_w = Inches(1.5)
    cell_w = (width - label_w) // n_cols
    cell_h = height // (n_rows + 1)

    flat = [v for row in values for v in row]
    v_min, v_max = min(flat), max(flat)

    for j, ch in enumerate(col_headers):
        hdr = slide.shapes.add_textbox(
            left + label_w + cell_w * j, top, cell_w, cell_h,
        )
        p = hdr.text_frame.paragraphs[0]
        p.text = ch
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.size = theme.font_size_caption
            run.font.bold = True
            run.font.color.rgb = theme.text_secondary
            run.font.name = theme.font_body

    for i, (rh, row) in enumerate(zip(row_headers, values)):
        row_top = top + cell_h * (i + 1)
        lbl = slide.shapes.add_textbox(left, row_top, label_w, cell_h)
        p = lbl.text_frame.paragraphs[0]
        p.text = rh
        for run in p.runs:
            run.font.size = theme.font_size_body
            run.font.color.rgb = theme.text_primary
            run.font.name = theme.font_body

        for j, val in enumerate(row):
            intensity = (val - v_min) / (v_max - v_min) if v_max > v_min else 0.5
            tint = int((1 - intensity) * 88)
            c = theme.primary
            tint_r = c[0] + (255 - c[0]) * tint // 100
            tint_g = c[1] + (255 - c[1]) * tint // 100
            tint_b = c[2] + (255 - c[2]) * tint // 100

            cell_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left + label_w + cell_w * j + Inches(0.03),
                row_top + Inches(0.03),
                cell_w - Inches(0.06),
                cell_h - Inches(0.06),
            )
            cell_shape.fill.solid()
            cell_shape.fill.fore_color.rgb = RGBColor(tint_r, tint_g, tint_b)
            cell_shape.line.fill.background()

            tf = cell_shape.text_frame
            p2 = tf.paragraphs[0]
            p2.alignment = PP_ALIGN.CENTER
            p2.text = str(int(val)) if val == int(val) else f"{val:.1f}"
            text_color = RGBColor(0xFF, 0xFF, 0xFF) if intensity >= 0.6 else theme.text_primary
            for run in p2.runs:
                run.font.size = theme.font_size_caption
                run.font.color.rgb = text_color
                run.font.name = theme.font_body
                run.font.bold = intensity >= 0.6


def add_benchmark_bar(slide, theme, items, left, top, width=None, height=None, unit=None):
    """横棒比較: 自社 (is_self=True) は primary 色、競合は neutral 色で強調
    items: [{"label": str, "value": float, "is_self": bool?}]
    """
    if width is None:
        width = theme.content_width
    if height is None:
        height = Inches(3.5)

    n = len(items)
    label_w = Inches(2.0)
    bar_area_w = width - label_w - Inches(0.8)
    row_h = height // n
    bar_h = row_h - Inches(0.12)
    max_val = max(item["value"] for item in items) if items else 1

    for i, item in enumerate(items):
        row_top = top + row_h * i
        is_self = item.get("is_self", False)

        lbl = slide.shapes.add_textbox(
            left, row_top + (row_h - Inches(0.3)) // 2, label_w, Inches(0.3),
        )
        p = lbl.text_frame.paragraphs[0]
        p.text = item["label"]
        for run in p.runs:
            run.font.size = theme.font_size_body
            run.font.color.rgb = theme.text_primary
            run.font.name = theme.font_body
            run.font.bold = is_self

        bar_w = max(int(bar_area_w * item["value"] / max_val), Inches(0.05))
        bar_top = row_top + (row_h - bar_h) // 2

        bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left + label_w, bar_top, bar_w, bar_h,
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = theme.primary if is_self else theme.neutral
        bar.line.fill.background()

        val_str = f"{item['value']}{unit}" if unit else str(item["value"])
        val_box = slide.shapes.add_textbox(
            left + label_w + bar_w + Inches(0.1),
            row_top + (row_h - Inches(0.3)) // 2,
            Inches(0.7), Inches(0.3),
        )
        p2 = val_box.text_frame.paragraphs[0]
        p2.text = val_str
        for run in p2.runs:
            run.font.size = theme.font_size_body
            run.font.color.rgb = theme.primary if is_self else theme.text_secondary
            run.font.name = theme.font_body
            run.font.bold = is_self
