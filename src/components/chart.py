from pptx.util import Inches, Pt, Emu
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from src.components._style import set_paragraph_text


# python-pptx チャートのプロット領域推定比率。左 15% は y 軸ラベル+目盛。
PLOT_AREA_LEFT_RATIO = 0.12
PLOT_AREA_RIGHT_RATIO = 0.98
PLOT_AREA_TOP_OFFSET = Inches(0.15)


def add_bar_chart(slide, theme, data, left, top, width=None, height=None,
                   horizontal=False, unit=None, annotations=None):
    """棒グラフ: 縦/横対応、データラベル付き、任意で注記オーバーレイ"""
    if width is None:
        width = Inches(8.0)
    if height is None:
        height = Inches(5.0)

    chart_data = CategoryChartData()
    chart_data.categories = data["labels"]
    for series in data["series"]:
        chart_data.add_series(series["name"], series["values"])

    chart_type = XL_CHART_TYPE.BAR_CLUSTERED if horizontal else XL_CHART_TYPE.COLUMN_CLUSTERED
    chart_frame = slide.shapes.add_chart(chart_type, left, top, width, height, chart_data)
    chart = chart_frame.chart

    _style_chart(chart, theme, data, unit)

    if annotations and not horizontal:
        _render_annotations(slide, theme, annotations, data["labels"], left, top, width, height)


def add_line_chart(slide, theme, data, left, top, width=None, height=None, unit=None, annotations=None):
    """折れ線グラフ: マーカー付き、任意で注記オーバーレイ"""
    if width is None:
        width = Inches(8.0)
    if height is None:
        height = Inches(5.0)

    chart_data = CategoryChartData()
    chart_data.categories = data["labels"]
    for series in data["series"]:
        chart_data.add_series(series["name"], series["values"])

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, left, top, width, height, chart_data
    )
    chart = chart_frame.chart

    _style_chart(chart, theme, data, unit)

    for series in chart.series:
        series.smooth = False
        series.format.line.width = Pt(2.5)

    if annotations:
        _render_annotations(slide, theme, annotations, data["labels"], left, top, width, height)


def add_stacked_bar_chart(slide, theme, data, left, top, width=None, height=None,
                           horizontal=False, unit=None, annotations=None):
    """積み上げ棒グラフ。複数系列を縦に積み上げる。"""
    if width is None:
        width = Inches(8.0)
    if height is None:
        height = Inches(5.0)

    chart_data = CategoryChartData()
    chart_data.categories = data["labels"]
    for series in data["series"]:
        chart_data.add_series(series["name"], series["values"])

    chart_type = XL_CHART_TYPE.BAR_STACKED if horizontal else XL_CHART_TYPE.COLUMN_STACKED
    chart_frame = slide.shapes.add_chart(chart_type, left, top, width, height, chart_data)
    chart = chart_frame.chart

    _style_chart(chart, theme, data, unit)

    # 積み上げでは個別ラベルが煩雑になるため値は非表示にして系列名のみ凡例で示す
    plot = chart.plots[0]
    plot.has_data_labels = False

    if annotations and not horizontal:
        _render_annotations(slide, theme, annotations, data["labels"], left, top, width, height)


def add_area_chart(slide, theme, data, left, top, width=None, height=None,
                    stacked=False, unit=None, annotations=None):
    """面グラフ。stacked=True で積み上げ面に切替。"""
    if width is None:
        width = Inches(8.0)
    if height is None:
        height = Inches(5.0)

    chart_data = CategoryChartData()
    chart_data.categories = data["labels"]
    for series in data["series"]:
        chart_data.add_series(series["name"], series["values"])

    chart_type = XL_CHART_TYPE.AREA_STACKED if stacked else XL_CHART_TYPE.AREA
    chart_frame = slide.shapes.add_chart(chart_type, left, top, width, height, chart_data)
    chart = chart_frame.chart

    _style_chart(chart, theme, data, unit)

    # 面の塗りに半透明感を出すため、線色を濃くして塗りはやや薄くする運用は theme 側で調整
    plot = chart.plots[0]
    plot.has_data_labels = False

    if annotations:
        _render_annotations(slide, theme, annotations, data["labels"], left, top, width, height)


def add_scatter_chart(slide, theme, data, left, top, width=None, height=None,
                       x_label=None, y_label=None):
    """散布図 (XY)。data は {series: [{name, points: [[x, y], ...]}]}"""
    from pptx.chart.data import XyChartData
    if width is None:
        width = Inches(8.0)
    if height is None:
        height = Inches(5.0)

    chart_data = XyChartData()
    for series in data["series"]:
        s = chart_data.add_series(series["name"])
        for x, y in series["points"]:
            s.add_data_point(x, y)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.XY_SCATTER, left, top, width, height, chart_data
    )
    chart = chart_frame.chart

    if len(data["series"]) > 1:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = theme.font_size_caption
        chart.legend.font.name = theme.font_body
    else:
        chart.has_legend = False

    value_axis = chart.value_axis
    value_axis.has_major_gridlines = True
    value_axis.major_gridlines.format.line.color.rgb = theme.border
    value_axis.major_gridlines.format.line.width = Pt(0.5)
    value_axis.tick_labels.font.size = theme.font_size_caption
    value_axis.tick_labels.font.name = theme.font_body
    if y_label:
        value_axis.has_title = True
        value_axis.axis_title.text_frame.text = y_label
        for p in value_axis.axis_title.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = theme.font_size_caption
                r.font.name = theme.font_body

    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = theme.font_size_caption
    category_axis.tick_labels.font.name = theme.font_body
    if x_label:
        category_axis.has_title = True
        category_axis.axis_title.text_frame.text = x_label
        for p in category_axis.axis_title.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = theme.font_size_caption
                r.font.name = theme.font_body

    for i, series in enumerate(chart.series):
        color_idx = i % len(theme.chart_colors)
        marker = series.marker
        marker.style = 8  # CIRCLE
        marker.size = 8
        marker.format.fill.solid()
        marker.format.fill.fore_color.rgb = theme.chart_colors[color_idx]
        marker.format.line.color.rgb = theme.chart_colors[color_idx]


def add_combo_chart(slide, theme, data, left, top, width=None, height=None, annotations=None):
    """棒+折れ線の複合グラフ。第2軸対応。

    data 形式:
        {
            "labels": [...],
            "bars":  [{"name": "売上", "values": [...], "unit": "億円"}],
            "lines": [{"name": "成長率", "values": [...], "unit": "%", "secondary_axis": true}]
        }
    """
    from copy import deepcopy
    from lxml import etree

    if width is None:
        width = Inches(8.0)
    if height is None:
        height = Inches(5.0)

    bars = data.get("bars", [])
    lines = data.get("lines", [])
    labels = data["labels"]

    # まず棒グラフとして全系列 (bars + lines) を入れて生成
    chart_data = CategoryChartData()
    chart_data.categories = labels
    for s in bars:
        chart_data.add_series(s["name"], s["values"])
    for s in lines:
        chart_data.add_series(s["name"], s["values"])

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart

    # 共通スタイル
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = theme.font_size_caption
    chart.legend.font.name = theme.font_body

    primary_axis = chart.value_axis
    primary_axis.has_major_gridlines = True
    primary_axis.major_gridlines.format.line.color.rgb = theme.border
    primary_axis.major_gridlines.format.line.width = Pt(0.5)
    primary_axis.tick_labels.font.size = theme.font_size_caption
    primary_axis.tick_labels.font.name = theme.font_body
    if bars and bars[0].get("unit"):
        primary_axis.tick_labels.number_format = f'#,##0"{bars[0]["unit"]}"'
        primary_axis.tick_labels.number_format_is_linked = False

    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = theme.font_size_caption
    category_axis.tick_labels.font.name = theme.font_body
    category_axis.format.line.color.rgb = theme.border

    # 系列の色付け
    for i, series in enumerate(chart.series):
        color_idx = i % len(theme.chart_colors)
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = theme.chart_colors[color_idx]

    # XML を直接操作して line 系列を折れ線化 + 第2軸対応
    nsmap = {"c": "http://schemas.openxmlformats.org/drawingml/2006/chart"}
    chart_xml = chart._chartSpace
    plot_area = chart_xml.find(".//c:plotArea", nsmap)
    bar_chart_el = plot_area.find("c:barChart", nsmap)
    bar_series_els = bar_chart_el.findall("c:ser", nsmap)

    n_bars = len(bars)
    line_series_els = bar_series_els[n_bars:]
    for el in line_series_els:
        bar_chart_el.remove(el)

    # 第2軸が必要かチェック
    use_secondary = any(s.get("secondary_axis") for s in lines)
    secondary_axis_id = "987654321" if use_secondary else None

    if lines:
        c = "{http://schemas.openxmlformats.org/drawingml/2006/chart}"
        line_chart_el = etree.SubElement(plot_area, f"{c}lineChart")
        etree.SubElement(line_chart_el, f"{c}grouping").set("val", "standard")
        etree.SubElement(line_chart_el, f"{c}varyColors").set("val", "0")

        for el in line_series_els:
            # marker と smooth を追加
            marker_el = etree.SubElement(el, f"{c}marker")
            etree.SubElement(marker_el, f"{c}symbol").set("val", "circle")
            etree.SubElement(marker_el, f"{c}size").set("val", "7")
            smooth_el = etree.SubElement(el, f"{c}smooth")
            smooth_el.set("val", "0")
            line_chart_el.append(el)

        etree.SubElement(line_chart_el, f"{c}marker").set("val", "1")

        # 軸IDを参照
        bar_axis_ids = bar_chart_el.findall("c:axId", nsmap)
        cat_axis_id = bar_axis_ids[0].get("val")
        primary_val_axis_id = bar_axis_ids[1].get("val")

        if use_secondary:
            ax_id_el1 = etree.SubElement(line_chart_el, f"{c}axId")
            ax_id_el1.set("val", cat_axis_id)
            ax_id_el2 = etree.SubElement(line_chart_el, f"{c}axId")
            ax_id_el2.set("val", secondary_axis_id)

            # 第2軸 (valAx) を追加
            sec_val_ax = etree.SubElement(plot_area, f"{c}valAx")
            etree.SubElement(sec_val_ax, f"{c}axId").set("val", secondary_axis_id)
            scaling = etree.SubElement(sec_val_ax, f"{c}scaling")
            etree.SubElement(scaling, f"{c}orientation").set("val", "minMax")
            etree.SubElement(sec_val_ax, f"{c}delete").set("val", "0")
            etree.SubElement(sec_val_ax, f"{c}axPos").set("val", "r")
            etree.SubElement(sec_val_ax, f"{c}crossAx").set("val", cat_axis_id)
            etree.SubElement(sec_val_ax, f"{c}crosses").set("val", "max")
            unit_str = lines[0].get("unit", "") if lines else ""
            if unit_str:
                fmt = etree.SubElement(sec_val_ax, f"{c}numFmt")
                fmt.set("formatCode", f'#,##0"{unit_str}"')
                fmt.set("sourceLinked", "0")
        else:
            ax_id_el1 = etree.SubElement(line_chart_el, f"{c}axId")
            ax_id_el1.set("val", cat_axis_id)
            ax_id_el2 = etree.SubElement(line_chart_el, f"{c}axId")
            ax_id_el2.set("val", primary_val_axis_id)

    if annotations:
        _render_annotations(slide, theme, annotations, labels, left, top, width, height)


def add_pie_chart(slide, theme, data, left, top, width=None, height=None):
    """円グラフ: ラベル+割合表示"""
    if width is None:
        width = Inches(6.0)
    if height is None:
        height = Inches(5.0)

    chart_data = CategoryChartData()
    chart_data.categories = data["labels"]
    chart_data.add_series("", data["values"])

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, left, top, width, height, chart_data
    )
    chart = chart_frame.chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = theme.font_size_caption
    chart.legend.font.name = theme.font_body

    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.show_percentage = True
    data_labels.show_category_name = False
    data_labels.show_value = False
    data_labels.font.size = theme.font_size_body
    data_labels.font.name = theme.font_body

    series = chart.series[0]
    for i, point in enumerate(series.points):
        color_idx = i % len(theme.chart_colors)
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = theme.chart_colors[color_idx]


def add_waterfall(slide, theme, data, left, top, width=None, height=None):
    """ウォーターフォールチャート: 積み上げ棒で増減を表現"""
    if width is None:
        width = Inches(8.0)
    if height is None:
        height = Inches(5.0)

    labels = data["labels"]
    values = data["values"]

    invisible = []
    visible = []
    running = 0
    for i, val in enumerate(values):
        if i == 0:
            invisible.append(0)
            visible.append(val)
            running = val
        elif i == len(values) - 1:
            invisible.append(0)
            visible.append(val)
        else:
            if val >= 0:
                invisible.append(running)
                visible.append(val)
                running += val
            else:
                running += val
                invisible.append(running)
                visible.append(abs(val))

    chart_data = CategoryChartData()
    chart_data.categories = labels
    chart_data.add_series("base", invisible)
    chart_data.add_series("value", visible)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart

    base_series = chart.series[0]
    base_series.format.fill.background()
    base_series.format.line.fill.background()

    value_series = chart.series[1]
    for i, val in enumerate(values):
        point = value_series.points[i]
        point.format.fill.solid()
        if i == 0 or i == len(values) - 1:
            point.format.fill.fore_color.rgb = theme.primary
        elif val >= 0:
            point.format.fill.fore_color.rgb = theme.success
        else:
            point.format.fill.fore_color.rgb = theme.danger

    chart.has_legend = False
    value_axis = chart.value_axis
    value_axis.has_major_gridlines = True
    value_axis.major_gridlines.format.line.color.rgb = theme.border
    value_axis.major_gridlines.format.line.width = Pt(0.5)
    value_axis.format.line.fill.background()
    value_axis.tick_labels.font.size = theme.font_size_caption
    value_axis.tick_labels.font.name = theme.font_body

    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = theme.font_size_caption
    category_axis.tick_labels.font.name = theme.font_body
    category_axis.format.line.color.rgb = theme.border


def _style_chart(chart, theme, data, unit=None):
    """チャート共通スタイリング"""
    if len(data["series"]) > 1:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = theme.font_size_caption
        chart.legend.font.name = theme.font_body
    else:
        chart.has_legend = False

    value_axis = chart.value_axis
    value_axis.has_major_gridlines = True
    value_axis.major_gridlines.format.line.color.rgb = theme.border
    value_axis.major_gridlines.format.line.width = Pt(0.5)
    value_axis.format.line.fill.background()
    value_axis.tick_labels.font.size = theme.font_size_caption
    value_axis.tick_labels.font.name = theme.font_body
    if unit:
        value_axis.tick_labels.number_format = f'#,##0"{unit}"'
        value_axis.tick_labels.number_format_is_linked = False

    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = theme.font_size_caption
    category_axis.tick_labels.font.name = theme.font_body
    category_axis.format.line.color.rgb = theme.border
    category_axis.has_major_gridlines = False

    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.show_value = True
    data_labels.show_category_name = False
    data_labels.font.size = theme.font_size_caption
    data_labels.font.name = theme.font_body
    data_labels.number_format = 'General'
    data_labels.number_format_is_linked = False

    for i, series in enumerate(chart.series):
        color_idx = i % len(theme.chart_colors)
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = theme.chart_colors[color_idx]


def _render_annotations(slide, theme, annotations, labels, chart_left, chart_top, chart_width, chart_height):
    """チャート上にピル型の注記を配置。
    annotations: [{"category": str|int, "text": str, "position": "top"|"bottom"}]
      - category は labels の要素か 0-based インデックス
      - position: "top"(デフォ) はチャート上端 / "bottom" は下端付近
    """
    plot_left = chart_left + int(chart_width * PLOT_AREA_LEFT_RATIO)
    plot_right = chart_left + int(chart_width * PLOT_AREA_RIGHT_RATIO)
    plot_span = plot_right - plot_left
    n = len(labels)

    pill_width = Inches(1.3)
    pill_height = Inches(0.32)

    for ann in annotations:
        cat = ann.get("category")
        text = ann.get("text", "")
        position = ann.get("position", "top")

        if isinstance(cat, int):
            idx = cat
        else:
            try:
                idx = labels.index(cat)
            except ValueError:
                continue
        if idx < 0 or idx >= n:
            continue

        cat_center_x = plot_left + plot_span * (idx + 0.5) // n
        pill_left = cat_center_x - pill_width // 2
        if position == "bottom":
            pill_top = chart_top + chart_height - Inches(0.9)
        else:
            pill_top = chart_top + PLOT_AREA_TOP_OFFSET

        pill = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            pill_left, pill_top, pill_width, pill_height,
        )
        pill.fill.solid()
        pill.fill.fore_color.rgb = theme.secondary
        pill.line.fill.background()

        tf = pill.text_frame
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_paragraph_text(
            p, text,
            size=theme.font_size_caption, color=RGBColor(0xFF, 0xFF, 0xFF),
            name=theme.font_body, bold=True,
        )
