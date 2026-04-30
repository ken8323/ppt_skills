from pptx import Presentation
from pptx.util import Inches

from src.themes import get_theme
from src.layouts import get_layout
from src.components.footer import add_page_footer
from src.components.source_note import render_source_note, has_source
from src.validator import validate_config
from src.linter import lint_config


BLANK_LAYOUT_INDEX = 6

# フッターを付与しないレイアウト (背景塗りつぶし系)
FOOTER_SKIP_LAYOUTS = {"cover", "section_divider"}

# source 注記をスキップするレイアウト (背景塗りつぶし系・トップ表紙等)
SOURCE_SKIP_LAYOUTS = {"cover", "section_divider"}


def _should_skip_footer(layout_name: str, data: dict) -> bool:
    if layout_name in FOOTER_SKIP_LAYOUTS:
        return True
    if layout_name == "closing" and data.get("type") == "thank_you":
        return True
    return False


def generate_pptx(
    config: dict,
    output_path: str,
    *,
    validate: bool = True,
    lint: bool = True,
    strict: bool = None,
) -> Presentation:
    """設定辞書からpptxを生成してファイル保存。

    config 構造:
        {
            "theme": "monotone" | "dark" | "colorful",
            "footer": "株式会社ABC | 社外秘",  # 任意。各ページ左下に表示
            "brand_name": "...",  # 任意。theme.brand_name を上書き
            "slides": [
                {"layout": "cover", "data": {...}},
                ...
            ],
        }

    Args:
        validate: True で schema.json + ビジネスルール検証を実行 (既定 True)。
                  検証失敗時は ConfigValidationError を送出し pptx は生成しない。
        lint:     True で警告レベルの静的チェック (オーバーフロー等) を実行し
                  stderr に警告を出力 (既定 True)。生成は継続。
        strict:   後方互換用。指定時は validate に同期する (deprecated)。
    """
    # 後方互換: strict=True/False が渡された場合は validate に反映
    if strict is not None:
        validate = strict

    if validate:
        validate_config(config)

    if lint:
        warnings = lint_config(config)
        if warnings:
            import sys
            print("[ppt_skills] Lint warnings:", file=sys.stderr)
            for w in warnings:
                print(f"  - {w}", file=sys.stderr)

    theme_name = config.get("theme", "monotone")
    theme = get_theme(theme_name)

    if "brand_name" in config:
        theme.brand_name = config["brand_name"]

    footer_text = config.get("footer", "")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[BLANK_LAYOUT_INDEX]

    slide_configs = config.get("slides", [])
    total = len(slide_configs)

    for idx, slide_cfg in enumerate(slide_configs, start=1):
        layout_name = slide_cfg.get("layout")
        data = slide_cfg.get("data", {})

        slide = prs.slides.add_slide(blank_layout)
        layout = get_layout(layout_name)
        layout.render(slide, theme, data)

        if layout_name not in SOURCE_SKIP_LAYOUTS and has_source(data):
            render_source_note(slide, theme, data)

        if not _should_skip_footer(layout_name, data):
            add_page_footer(slide, theme, idx, total, footer_text=footer_text)

        notes_text = slide_cfg.get("notes", "")
        if notes_text:
            notes_slide = slide.notes_slide
            tf = notes_slide.notes_text_frame
            tf.text = notes_text

    prs.save(output_path)
    return prs
