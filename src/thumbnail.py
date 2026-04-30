"""PPTX サムネイル生成 (PNG)。

優先順位:
  1. LibreOffice headless (クロスプラットフォーム / 高品質)
  2. macOS qlmanage (ファイル全体の1枚サムネイルのみ)
  3. Pillow によるミニマルサムネイル (タイトル + テーマ背景色)
"""
from __future__ import annotations

import os
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Optional


class ThumbnailError(RuntimeError):
    pass


def generate_thumbnails(
    pptx_path: str | Path,
    output_dir: str | Path,
    *,
    dpi: int = 150,
    slide_index: Optional[int] = None,
) -> list[Path]:
    """PPTX の各スライドを PNG に変換して output_dir に保存。

    Args:
        pptx_path:   入力 .pptx ファイルのパス。
        output_dir:  出力ディレクトリ (存在しなければ作成)。
        dpi:         解像度 (LibreOffice 使用時に有効)。
        slide_index: None = 全スライド。0-based int で 1 枚だけ指定可。

    Returns:
        生成した PNG ファイルのパスリスト。
    """
    pptx_path = Path(pptx_path)
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    if shutil.which("libreoffice") or shutil.which("soffice"):
        return _via_libreoffice(pptx_path, output_dir, dpi, slide_index)
    if _is_macos() and shutil.which("qlmanage"):
        return _via_qlmanage(pptx_path, output_dir, slide_index)
    try:
        import PIL  # noqa: F401
        return _via_pillow(pptx_path, output_dir, slide_index)
    except ImportError:
        pass

    raise ThumbnailError(
        "サムネイル生成には LibreOffice か Pillow が必要です。\n"
        "  LibreOffice: https://www.libreoffice.org/download/download/\n"
        "  Pillow:      pip install Pillow"
    )


# ── LibreOffice ────────────────────────────────────────────────────────────────

def _via_libreoffice(
    pptx_path: Path, output_dir: Path, dpi: int, slide_index: Optional[int]
) -> list[Path]:
    cmd = shutil.which("libreoffice") or shutil.which("soffice")
    with tempfile.TemporaryDirectory() as tmp:
        subprocess.run(
            [cmd, "--headless", "--convert-to", "png",
             "--outdir", tmp, str(pptx_path)],
            check=True, capture_output=True,
        )
        tmp_pngs = sorted(Path(tmp).glob("*.png"))
        results: list[Path] = []
        for i, src in enumerate(tmp_pngs):
            if slide_index is not None and i != slide_index:
                continue
            dst = output_dir / f"slide_{i + 1:03d}.png"
            shutil.copy2(src, dst)
            results.append(dst)
    return results


# ── qlmanage (macOS) ───────────────────────────────────────────────────────────

def _is_macos() -> bool:
    import platform
    return platform.system() == "Darwin"


def _via_qlmanage(
    pptx_path: Path, output_dir: Path, slide_index: Optional[int]
) -> list[Path]:
    subprocess.run(
        ["qlmanage", "-t", "-s", "1920", "-o", str(output_dir), str(pptx_path)],
        check=True, capture_output=True,
    )
    pngs = sorted(output_dir.glob(f"{pptx_path.name}*.png"))
    if slide_index is not None:
        pngs = pngs[slide_index:slide_index + 1]
    return pngs


# ── Pillow フォールバック ──────────────────────────────────────────────────────

def _hex_to_rgb(hex_str: str) -> tuple[int, int, int]:
    h = hex_str.lstrip("#")
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)


def _via_pillow(
    pptx_path: Path, output_dir: Path, slide_index: Optional[int]
) -> list[Path]:
    from PIL import Image, ImageDraw, ImageFont
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation(str(pptx_path))
    W, H = 1280, 720
    results: list[Path] = []

    for i, slide in enumerate(prs.slides):
        if slide_index is not None and i != slide_index:
            continue

        # 背景色の推定 (塗りつぶし Shape があれば使用)
        bg_color = (240, 240, 240)
        for shape in slide.shapes:
            try:
                fill = shape.fill
                if fill.type is not None and fill.fore_color.type is not None:
                    rgb = fill.fore_color.rgb
                    bg_color = (rgb[0], rgb[1], rgb[2])
                    break
            except Exception:
                pass

        img = Image.new("RGB", (W, H), color=bg_color)
        draw = ImageDraw.Draw(img)

        # タイトルテキストを描画
        title_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                title_text = shape.text.strip().split("\n")[0]
                break

        font_size = 48
        try:
            font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", font_size)
        except Exception:
            font = ImageFont.load_default()

        text_color = (50, 50, 50) if sum(bg_color) > 400 else (220, 220, 220)
        draw.text((W // 2, H // 2), title_text, fill=text_color, font=font, anchor="mm")

        slide_num_text = f"{i + 1} / {len(prs.slides)}"
        draw.text((W - 40, H - 30), slide_num_text, fill=text_color,
                  font=ImageFont.load_default(), anchor="rm")

        dst = output_dir / f"slide_{i + 1:03d}.png"
        img.save(dst, "PNG")
        results.append(dst)

    return results
