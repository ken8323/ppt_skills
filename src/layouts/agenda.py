from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

from src.components.text import add_title


class AgendaLayout:
    def render(self, slide, theme, data):
        """アジェンダ: タイトル + 番号付きリスト"""
        add_title(slide, theme, "Agenda", theme.margin_left, theme.margin_top)

        items = data.get("items", [])
        highlight = data.get("highlight", None)

        item_top = theme.content_area_top + Inches(0.2)
        item_height = Inches(0.7)

        for i, item_text in enumerate(items):
            y = item_top + item_height * i
            is_active = (highlight is not None and i == highlight)

            num_width = Inches(0.8)
            num_box = slide.shapes.add_textbox(
                theme.margin_left + Inches(0.5), y, num_width, item_height,
            )
            num_box.text_frame.paragraphs[0].text = f"{i + 1:02d}"
            num_box.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
            for run in num_box.text_frame.paragraphs[0].runs:
                run.font.size = Pt(28)
                run.font.bold = True
                run.font.name = theme.font_title
                if is_active:
                    run.font.color.rgb = theme.secondary
                else:
                    run.font.color.rgb = theme.primary

            text_box = slide.shapes.add_textbox(
                theme.margin_left + Inches(1.6), y + Inches(0.1),
                Inches(8.0), item_height,
            )
            text_box.text_frame.paragraphs[0].text = item_text
            for run in text_box.text_frame.paragraphs[0].runs:
                run.font.size = Pt(20)
                run.font.name = theme.font_body
                if is_active:
                    run.font.color.rgb = theme.text_primary
                    run.font.bold = True
                else:
                    run.font.color.rgb = theme.text_secondary

            if i < len(items) - 1:
                sep = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    theme.margin_left + Inches(1.6), y + item_height - Pt(1),
                    Inches(8.0), Pt(1),
                )
                sep.fill.solid()
                sep.fill.fore_color.rgb = theme.border
                sep.line.fill.background()
