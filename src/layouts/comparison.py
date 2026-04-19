from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

from src.components.text import add_title
from src.layouts.content import ContentLayout


class ComparisonLayout:
    def render(self, slide, theme, data):
        """比較ページ: 左右2分割、中央に区切り線"""
        title = data.get("title", "")
        left_title = data.get("left_title", "")
        right_title = data.get("right_title", "")
        left_components = data.get("left_components", [])
        right_components = data.get("right_components", [])

        if title:
            add_title(slide, theme, title, theme.margin_left, theme.margin_top)

        content_top = theme.content_area_top
        half_width = (theme.content_width - Inches(0.6)) // 2

        if left_title:
            lt_box = slide.shapes.add_textbox(
                theme.margin_left, content_top, half_width, Inches(0.5),
            )
            lt_box.text_frame.paragraphs[0].text = left_title
            lt_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            for run in lt_box.text_frame.paragraphs[0].runs:
                run.font.size = Pt(18)
                run.font.bold = True
                run.font.color.rgb = theme.primary
                run.font.name = theme.font_title

        right_left = theme.margin_left + half_width + Inches(0.6)
        if right_title:
            rt_box = slide.shapes.add_textbox(
                right_left, content_top, half_width, Inches(0.5),
            )
            rt_box.text_frame.paragraphs[0].text = right_title
            rt_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            for run in rt_box.text_frame.paragraphs[0].runs:
                run.font.size = Pt(18)
                run.font.bold = True
                run.font.color.rgb = theme.secondary
                run.font.name = theme.font_title

        divider_x = theme.margin_left + half_width + Inches(0.25)
        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, divider_x, content_top, Pt(2), theme.content_height,
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = theme.border
        divider.line.fill.background()

        comp_top = content_top + Inches(0.7)
        content_layout = ContentLayout()
        content_layout._render_components(
            slide, theme, left_components,
            theme.margin_left, comp_top, half_width,
        )
        content_layout._render_components(
            slide, theme, right_components,
            right_left, comp_top, half_width,
        )
