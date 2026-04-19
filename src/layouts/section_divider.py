from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN


class SectionDividerLayout:
    def render(self, slide, theme, data):
        """セクション区切り: primary色背景、中央にセクション番号+名前"""
        sw = theme.slide_width
        sh = theme.slide_height

        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, sh)
        bg.fill.solid()
        bg.fill.fore_color.rgb = theme.primary
        bg.line.fill.background()

        section_number = data.get("section_number", "")
        section_title = data.get("section_title", "")

        center_width = Inches(8.0)
        center_left = (sw - center_width) // 2

        if section_number != "":
            num_box = slide.shapes.add_textbox(
                center_left, sh // 2 - Inches(1.5), center_width, Inches(1.0),
            )
            num_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            if isinstance(section_number, int):
                num_box.text_frame.paragraphs[0].text = f"{section_number:02d}"
            else:
                num_box.text_frame.paragraphs[0].text = str(section_number)
            for run in num_box.text_frame.paragraphs[0].runs:
                run.font.size = Pt(60)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                run.font.name = theme.font_title

        title_box = slide.shapes.add_textbox(
            center_left, sh // 2 - Inches(0.3), center_width, Inches(0.8),
        )
        title_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_box.text_frame.paragraphs[0].text = section_title
        for run in title_box.text_frame.paragraphs[0].runs:
            run.font.size = Pt(32)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            run.font.name = theme.font_title

        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            (sw - Inches(3.0)) // 2, sh // 2 + Inches(0.7),
            Inches(3.0), Pt(3),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = theme.secondary
        line.line.fill.background()
