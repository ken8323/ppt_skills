from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from src.components._style import set_paragraph_text


class CoverLayout:
    def render(self, slide, theme, data):
        """表紙: 左パネル(primary色) + 右側タイトル・クライアント情報"""
        sw = theme.slide_width
        sh = theme.slide_height

        # White background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, sh)
        bg.fill.solid()
        bg.fill.fore_color.rgb = theme.background
        bg.line.fill.background()

        # Left primary panel (~38% of width)
        panel_w = int(sw * 0.38)
        panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, panel_w, sh)
        panel.fill.solid()
        panel.fill.fore_color.rgb = theme.primary
        panel.line.fill.background()

        # Thin secondary accent strip at right edge of panel
        strip_w = Inches(0.07)
        strip = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, panel_w - strip_w, 0, strip_w, sh
        )
        strip.fill.solid()
        strip.fill.fore_color.rgb = theme.secondary
        strip.line.fill.background()

        # Brand name on left panel (top-left)
        brand_name = getattr(theme, "brand_name", "") or data.get("brand_name", "")
        if brand_name:
            brand_box = slide.shapes.add_textbox(
                Inches(0.35), Inches(0.35), panel_w - Inches(0.5), Inches(0.5)
            )
            set_paragraph_text(
                brand_box.text_frame.paragraphs[0], brand_name,
                size=Pt(13), color=RGBColor(0xFF, 0xFF, 0xFF),
                name=theme.font_title, bold=True,
            )

        # Right content area
        content_left = panel_w + Inches(0.6)
        content_w = sw - content_left - Inches(0.5)

        # Main title — vertically centered slightly above midpoint
        title_top = sh // 2 - Inches(1.5)
        title_box = slide.shapes.add_textbox(content_left, title_top, content_w, Inches(1.7))
        title_box.text_frame.word_wrap = True
        set_paragraph_text(
            title_box.text_frame.paragraphs[0], data.get("title", ""),
            size=Pt(30), color=theme.text_primary,
            name=theme.font_title, bold=True,
        )

        # Subtitle
        subtitle = data.get("subtitle", "")
        if subtitle:
            sub_box = slide.shapes.add_textbox(
                content_left, title_top + Inches(1.75), content_w, Inches(0.65)
            )
            sub_box.text_frame.word_wrap = True
            set_paragraph_text(
                sub_box.text_frame.paragraphs[0], subtitle,
                size=Pt(15), color=theme.text_secondary,
                name=theme.font_body,
            )

        # Separator line
        sep_top = sh // 2 + Inches(0.65)
        sep = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, content_left, sep_top, content_w, Pt(1.5)
        )
        sep.fill.solid()
        sep.fill.fore_color.rgb = theme.border
        sep.line.fill.background()

        # Client name
        client = data.get("client", "")
        if client:
            cl_box = slide.shapes.add_textbox(
                content_left, sep_top + Inches(0.15), content_w, Inches(0.4)
            )
            set_paragraph_text(
                cl_box.text_frame.paragraphs[0], client,
                size=Pt(13), color=theme.text_secondary,
                name=theme.font_body, bold=True,
            )

        # Date
        date = data.get("date", "")
        if date:
            dt_box = slide.shapes.add_textbox(
                content_left, sep_top + Inches(0.65), content_w, Inches(0.35)
            )
            set_paragraph_text(
                dt_box.text_frame.paragraphs[0], date,
                size=Pt(12), color=theme.text_secondary,
                name=theme.font_body,
            )
