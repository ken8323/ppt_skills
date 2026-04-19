from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

from src.components.text import add_title, add_bullets


class ClosingLayout:
    def render(self, slide, theme, data):
        """まとめ/Next Steps or Thank You"""
        close_type = data.get("type", "summary")

        if close_type == "thank_you":
            self._render_thank_you(slide, theme, data)
        else:
            self._render_summary(slide, theme, data)

    def _render_summary(self, slide, theme, data):
        summary = data.get("summary", [])
        next_steps = data.get("next_steps", [])

        add_title(slide, theme, "Summary", theme.margin_left, theme.margin_top)

        if summary:
            add_bullets(
                slide, theme, summary,
                theme.margin_left, theme.content_area_top,
                width=theme.content_width,
                height=Inches(0.4 * len(summary)),
            )

        if next_steps:
            ns_top = theme.content_area_top + Inches(0.4 * len(summary)) + Inches(0.8)

            ns_header = slide.shapes.add_textbox(
                theme.margin_left, ns_top - Inches(0.5),
                Inches(3.0), Inches(0.4),
            )
            ns_header.text_frame.paragraphs[0].text = "Next Steps"
            for run in ns_header.text_frame.paragraphs[0].runs:
                run.font.size = Pt(18)
                run.font.bold = True
                run.font.color.rgb = theme.primary
                run.font.name = theme.font_title

            for i, step in enumerate(next_steps):
                step_top = ns_top + Inches(0.5 * i)

                num_size = Inches(0.35)
                num_shape = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    theme.margin_left, step_top, num_size, num_size,
                )
                num_shape.fill.solid()
                num_shape.fill.fore_color.rgb = theme.primary
                num_shape.line.fill.background()
                num_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                num_shape.text_frame.paragraphs[0].text = str(i + 1)
                for run in num_shape.text_frame.paragraphs[0].runs:
                    run.font.size = Pt(12)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    run.font.name = theme.font_title

                step_box = slide.shapes.add_textbox(
                    theme.margin_left + Inches(0.6), step_top,
                    theme.content_width - Inches(0.6), Inches(0.4),
                )
                step_box.text_frame.paragraphs[0].text = step
                for run in step_box.text_frame.paragraphs[0].runs:
                    run.font.size = theme.font_size_body
                    run.font.color.rgb = theme.text_primary
                    run.font.name = theme.font_body

    def _render_thank_you(self, slide, theme, data):
        sw = theme.slide_width
        sh = theme.slide_height

        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, sh)
        bg.fill.solid()
        bg.fill.fore_color.rgb = theme.primary
        bg.line.fill.background()

        ty_width = Inches(8.0)
        ty_left = (sw - ty_width) // 2
        ty_box = slide.shapes.add_textbox(ty_left, sh // 2 - Inches(1.0), ty_width, Inches(1.0))
        ty_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        ty_box.text_frame.paragraphs[0].text = "Thank You"
        for run in ty_box.text_frame.paragraphs[0].runs:
            run.font.size = Pt(48)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            run.font.name = theme.font_title

        contact = data.get("contact", "")
        if contact:
            ct_box = slide.shapes.add_textbox(ty_left, sh // 2 + Inches(0.3), ty_width, Inches(0.5))
            ct_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            ct_box.text_frame.paragraphs[0].text = contact
            for run in ct_box.text_frame.paragraphs[0].runs:
                run.font.size = Pt(16)
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                run.font.name = theme.font_body
