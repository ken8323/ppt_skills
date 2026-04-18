import pytest
from pptx import Presentation
from pptx.util import Inches

from src.components.text import add_title, add_subtitle, add_bullets, add_callout, add_footnote
from src.themes.monotone import MonotoneTheme


@pytest.fixture
def theme():
    return MonotoneTheme()


@pytest.fixture
def slide():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)


class TestAddTitle:
    def test_adds_textbox(self, slide, theme):
        add_title(slide, theme, "テストタイトル", Inches(0.6), Inches(0.4))
        assert len(slide.shapes) >= 1

    def test_text_content(self, slide, theme):
        add_title(slide, theme, "テストタイトル", Inches(0.6), Inches(0.4))
        text = slide.shapes[0].text_frame.paragraphs[0].text
        assert text == "テストタイトル"

    def test_adds_separator_line(self, slide, theme):
        add_title(slide, theme, "テストタイトル", Inches(0.6), Inches(0.4))
        assert len(slide.shapes) >= 2


class TestAddSubtitle:
    def test_adds_textbox(self, slide, theme):
        add_subtitle(slide, theme, "サブタイトル", Inches(0.6), Inches(1.0))
        assert len(slide.shapes) >= 1

    def test_text_content(self, slide, theme):
        add_subtitle(slide, theme, "サブタイトル", Inches(0.6), Inches(1.0))
        text = slide.shapes[0].text_frame.paragraphs[0].text
        assert text == "サブタイトル"


class TestAddBullets:
    def test_adds_textbox(self, slide, theme):
        add_bullets(slide, theme, ["項目1", "項目2", "項目3"], Inches(0.6), Inches(1.5))
        assert len(slide.shapes) >= 1

    def test_bullet_count(self, slide, theme):
        items = ["項目1", "項目2", "項目3"]
        add_bullets(slide, theme, items, Inches(0.6), Inches(1.5))
        paragraphs = slide.shapes[0].text_frame.paragraphs
        texts = [p.text for p in paragraphs if p.text.strip()]
        assert len(texts) == 3

    def test_nested_bullets(self, slide, theme):
        items = ["項目1", {"text": "サブ項目", "level": 1}]
        add_bullets(slide, theme, items, Inches(0.6), Inches(1.5))
        assert len(slide.shapes) >= 1


class TestAddCallout:
    def test_adds_shape(self, slide, theme):
        add_callout(slide, theme, "重要なメッセージ", Inches(0.6), Inches(2.0))
        assert len(slide.shapes) >= 1


class TestAddFootnote:
    def test_adds_textbox(self, slide, theme):
        add_footnote(slide, theme, "出典: 調査レポート 2026", Inches(0.6))
        assert len(slide.shapes) >= 1
