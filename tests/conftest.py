import pytest
from pptx import Presentation
from pptx.util import Inches

from src.themes.monotone import MonotoneTheme
from src.themes.dark import DarkTheme
from src.themes.colorful import ColorfulTheme


@pytest.fixture
def prs():
    p = Presentation()
    p.slide_width = Inches(13.333)
    p.slide_height = Inches(7.5)
    return p


@pytest.fixture
def blank_slide(prs):
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)


@pytest.fixture
def monotone_theme():
    return MonotoneTheme()


@pytest.fixture
def dark_theme():
    return DarkTheme()


@pytest.fixture
def colorful_theme():
    return ColorfulTheme()
