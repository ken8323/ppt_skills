import pytest
from pptx import Presentation
from pptx.util import Inches

from src.components.table import add_table
from src.themes.monotone import MonotoneTheme


@pytest.fixture
def theme():
    return MonotoneTheme()


@pytest.fixture
def slide():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)


class TestAddTable:
    def test_adds_table_shape(self, slide, theme):
        headers = ["項目", "2024", "2025"]
        rows = [["売上", "100", "150"], ["利益", "20", "35"]]
        add_table(slide, theme, headers, rows, Inches(0.6), Inches(1.5))
        assert len(slide.shapes) >= 1

    def test_correct_dimensions(self, slide, theme):
        headers = ["項目", "2024", "2025"]
        rows = [["売上", "100", "150"], ["利益", "20", "35"]]
        add_table(slide, theme, headers, rows, Inches(0.6), Inches(1.5))
        table = slide.shapes[0].table
        assert len(table.rows) == 3
        assert len(table.columns) == 3

    def test_header_text(self, slide, theme):
        headers = ["項目", "2024", "2025"]
        rows = [["売上", "100", "150"]]
        add_table(slide, theme, headers, rows, Inches(0.6), Inches(1.5))
        table = slide.shapes[0].table
        assert table.cell(0, 0).text == "項目"
        assert table.cell(0, 1).text == "2024"

    def test_data_text(self, slide, theme):
        headers = ["項目", "値"]
        rows = [["売上", "100"]]
        add_table(slide, theme, headers, rows, Inches(0.6), Inches(1.5))
        table = slide.shapes[0].table
        assert table.cell(1, 0).text == "売上"
        assert table.cell(1, 1).text == "100"
