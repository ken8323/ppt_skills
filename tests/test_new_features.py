"""Tier 1 + 2.1/2.3/2.2/2.4/2.5/2.6/2.7/2.8 で追加した機能の個別テスト。"""
import pytest
from pptx.util import Inches
from pptx.dml.color import RGBColor

from src.themes.base import Grid
from src.themes.monotone import MonotoneTheme
from src.components.text import add_title, add_callout
from src.components.chart import add_bar_chart, add_line_chart
from src.components.icon import add_kpi_card
from src.components.footer import add_page_footer
from src.components.table import add_table
from src.components.shape import (
    add_matrix_2x2, add_pyramid, add_process_flow,
    add_pillars, add_swot, add_heatmap, add_benchmark_bar,
)
from src.components.timeline import add_timeline, add_gantt
from src.layouts.chart_page import ChartPageLayout
from src.layouts.content import ContentLayout


class TestGrid:
    def test_grid_returns_twelve_columns(self, monotone_theme):
        grid = monotone_theme.grid
        assert grid.TOTAL_COLS == 12

    def test_col_x_increases(self, monotone_theme):
        grid = monotone_theme.grid
        assert grid.col_x(3) > grid.col_x(0)

    def test_span_width_scales(self, monotone_theme):
        grid = monotone_theme.grid
        assert grid.span_width(6) > grid.span_width(3)

    def test_cell_returns_tuple(self, monotone_theme):
        grid = monotone_theme.grid
        left, top, width, height = grid.cell(0, Inches(1), span_cols=4, height=Inches(2))
        assert left == monotone_theme.margin_left
        assert width > 0


class TestSemanticPalette:
    def test_monotone_has_success(self, monotone_theme):
        assert isinstance(monotone_theme.success, RGBColor)
        assert isinstance(monotone_theme.warning, RGBColor)
        assert isinstance(monotone_theme.danger, RGBColor)

    def test_dark_has_success(self, dark_theme):
        assert isinstance(dark_theme.success, RGBColor)


class TestFontHierarchy:
    def test_all_six_sizes_exist(self, monotone_theme):
        for attr in ["font_size_h1", "font_size_h2", "font_size_h3",
                     "font_size_body", "font_size_caption", "font_size_footnote"]:
            assert hasattr(monotone_theme, attr)

    def test_legacy_aliases_work(self, monotone_theme):
        assert monotone_theme.font_size_title == monotone_theme.font_size_h1
        assert monotone_theme.font_size_subtitle == monotone_theme.font_size_h2


class TestTitleSubhead:
    def test_title_with_subtitle_adds_extra_shape(self, blank_slide, monotone_theme):
        before = len(blank_slide.shapes)
        add_title(blank_slide, monotone_theme, "タイトル", Inches(0.5), Inches(0.5))
        mid = len(blank_slide.shapes)
        add_title(blank_slide, monotone_theme, "タイトル2", Inches(0.5), Inches(2.0),
                  subtitle="サブヘッド一行")
        after = len(blank_slide.shapes)
        assert (after - mid) > (mid - before), "subtitle 指定時は追加のテキストボックスが出る"


class TestCalloutVariants:
    @pytest.mark.parametrize("variant", ["info", "success", "warning", "danger"])
    def test_all_variants_render(self, blank_slide, monotone_theme, variant):
        before = len(blank_slide.shapes)
        add_callout(blank_slide, monotone_theme, "テキスト",
                    Inches(0.5), Inches(0.5), variant=variant)
        assert len(blank_slide.shapes) > before


class TestFooter:
    def test_page_footer_adds_shapes(self, blank_slide, monotone_theme):
        before = len(blank_slide.shapes)
        add_page_footer(blank_slide, monotone_theme, 3, 20, footer_text="社外秘")
        assert len(blank_slide.shapes) > before


class TestChartAnnotations:
    def test_bar_with_annotations(self, blank_slide, monotone_theme):
        data = {
            "labels": ["2023", "2024", "2025"],
            "series": [{"name": "売上", "values": [100, 150, 220]}],
        }
        before = len(blank_slide.shapes)
        add_bar_chart(blank_slide, monotone_theme, data,
                      Inches(0.5), Inches(0.5), width=Inches(8), height=Inches(4),
                      annotations=[{"category": "2025", "text": "+48%"}])
        assert len(blank_slide.shapes) > before + 1  # chart + 注記ピル

    def test_bar_annotation_by_index(self, blank_slide, monotone_theme):
        data = {
            "labels": ["Q1", "Q2", "Q3", "Q4"],
            "series": [{"name": "実績", "values": [10, 20, 30, 40]}],
        }
        add_bar_chart(blank_slide, monotone_theme, data,
                      Inches(0.5), Inches(0.5), width=Inches(8), height=Inches(4),
                      annotations=[{"category": 3, "text": "最大"}])

    def test_line_with_annotations(self, blank_slide, monotone_theme):
        data = {
            "labels": ["1月", "2月", "3月"],
            "series": [{"name": "件数", "values": [10, 15, 25]}],
        }
        add_line_chart(blank_slide, monotone_theme, data,
                       Inches(0.5), Inches(0.5), width=Inches(8), height=Inches(4),
                       annotations=[{"category": "3月", "text": "+150%", "position": "bottom"}])


class TestKPIDelta:
    def test_kpi_with_delta_up(self, blank_slide, monotone_theme):
        before = len(blank_slide.shapes)
        add_kpi_card(blank_slide, monotone_theme, "120", "億円", "売上",
                     Inches(0.5), Inches(0.5),
                     delta="+12%", delta_direction="up")
        assert len(blank_slide.shapes) >= before + 2  # card + delta badge

    def test_kpi_without_delta_skips_badge(self, blank_slide, monotone_theme):
        before = len(blank_slide.shapes)
        add_kpi_card(blank_slide, monotone_theme, "120", "億円", "売上",
                     Inches(0.5), Inches(0.5))
        assert len(blank_slide.shapes) == before + 1  # card のみ

    @pytest.mark.parametrize("direction", ["up", "down", "flat"])
    def test_kpi_all_directions(self, blank_slide, monotone_theme, direction):
        add_kpi_card(blank_slide, monotone_theme, "100", "%", "達成率",
                     Inches(0.5), Inches(0.5),
                     delta="+5", delta_direction=direction)


class TestChartPageSource:
    def test_chart_page_with_source(self, blank_slide, monotone_theme):
        data = {
            "title": "売上推移",
            "chart": {
                "type": "bar",
                "data": {
                    "labels": ["2023", "2024"],
                    "series": [{"name": "売上", "values": [100, 150]}],
                },
            },
            "source": "社内データ (2025年4月時点)",
        }
        before = len(blank_slide.shapes)
        ChartPageLayout().render(blank_slide, monotone_theme, data)
        assert len(blank_slide.shapes) > before


class TestTableHighlight:
    def test_highlight_rows_renders(self, blank_slide, monotone_theme):
        before = len(blank_slide.shapes)
        add_table(
            blank_slide, monotone_theme,
            ["A", "B"], [["1", "2"], ["3", "4"], ["5", "6"]],
            Inches(0.5), Inches(0.5),
            highlight_rows=[0, 2],
        )
        assert len(blank_slide.shapes) > before

    def test_highlight_cells_renders(self, blank_slide, monotone_theme):
        before = len(blank_slide.shapes)
        add_table(
            blank_slide, monotone_theme,
            ["X", "Y"], [["a", "b"], ["c", "d"]],
            Inches(0.5), Inches(0.5),
            highlight_cells=[{"row": 0, "col": 1}],
        )
        assert len(blank_slide.shapes) > before

    def test_no_highlight_unchanged(self, blank_slide, monotone_theme):
        before = len(blank_slide.shapes)
        add_table(
            blank_slide, monotone_theme,
            ["X"], [["a"], ["b"]],
            Inches(0.5), Inches(0.5),
        )
        assert len(blank_slide.shapes) > before


class TestMatrixRecommended:
    def test_recommended_quadrant_renders(self, blank_slide, monotone_theme):
        before = len(blank_slide.shapes)
        add_matrix_2x2(
            blank_slide, monotone_theme,
            x_axis="X軸", y_axis="Y軸",
            quadrants=["A", "B", "C", "D"],
            left=Inches(1), top=Inches(1),
            recommended_quadrant=1,
        )
        assert len(blank_slide.shapes) > before


class TestPyramidNotes:
    def test_pyramid_with_dict_levels(self, blank_slide, monotone_theme):
        before = len(blank_slide.shapes)
        add_pyramid(
            blank_slide, monotone_theme,
            levels=[
                {"text": "戦略", "note": "最上位目標"},
                {"text": "戦術", "note": "実行手段"},
                {"text": "施策", "note": "具体アクション"},
            ],
            left=Inches(1), top=Inches(1),
        )
        assert len(blank_slide.shapes) > before

    def test_pyramid_mixed_levels(self, blank_slide, monotone_theme):
        add_pyramid(
            blank_slide, monotone_theme,
            levels=["トップ", {"text": "ミドル", "note": "補足"}, "ボトム"],
            left=Inches(1), top=Inches(1),
        )


class TestProcessFlowChevron:
    def test_chevron_style_renders(self, blank_slide, monotone_theme):
        before = len(blank_slide.shapes)
        add_process_flow(
            blank_slide, monotone_theme,
            steps=["計画", "実行", "評価", "改善"],
            left=Inches(0.5), top=Inches(1),
            width=Inches(12),
            style="chevron",
        )
        assert len(blank_slide.shapes) > before

    def test_arrow_style_still_works(self, blank_slide, monotone_theme):
        before = len(blank_slide.shapes)
        add_process_flow(
            blank_slide, monotone_theme,
            steps=["STEP1", "STEP2", "STEP3"],
            left=Inches(0.5), top=Inches(1),
            style="arrow",
        )
        assert len(blank_slide.shapes) > before


class TestTimelineToday:
    def test_today_marker_adds_shapes(self, blank_slide, monotone_theme):
        before = len(blank_slide.shapes)
        add_timeline(
            blank_slide, monotone_theme,
            milestones=[
                {"date": "2026Q2", "label": "設計"},
                {"date": "2026Q4", "label": "開発"},
                {"date": "2027Q2", "label": "展開"},
            ],
            left=Inches(0.5), top=Inches(2),
            today="2026Q3",
        )
        after = len(blank_slide.shapes)
        assert after > before + 2  # ドット+ラベル以外にマーカーが追加されている

    def test_no_today_no_marker(self, blank_slide, monotone_theme):
        before = len(blank_slide.shapes)
        add_timeline(
            blank_slide, monotone_theme,
            milestones=[{"date": "2026Q1", "label": "開始"}, {"date": "2026Q3", "label": "完了"}],
            left=Inches(0.5), top=Inches(2),
        )
        assert len(blank_slide.shapes) > before


class TestGanttProgress:
    def test_progress_adds_extra_bar(self, blank_slide, monotone_theme):
        before = len(blank_slide.shapes)
        add_gantt(
            blank_slide, monotone_theme,
            tasks=[
                {"name": "設計", "start": 0, "duration": 2, "progress": 1.0},
                {"name": "開発", "start": 2, "duration": 3, "progress": 0.5},
                {"name": "テスト", "start": 5, "duration": 2},
            ],
            phases=["1月", "2月", "3月", "4月", "5月", "6月", "7月"],
            left=Inches(0.5), top=Inches(1),
        )
        assert len(blank_slide.shapes) > before

    def test_zero_progress_skips_bar(self, blank_slide, monotone_theme):
        before = len(blank_slide.shapes)
        add_gantt(
            blank_slide, monotone_theme,
            tasks=[{"name": "未着手", "start": 0, "duration": 2, "progress": 0.0}],
            phases=["A", "B"],
            left=Inches(0.5), top=Inches(1),
        )
        count_no_progress = len(blank_slide.shapes) - before

        before2 = len(blank_slide.shapes)
        add_gantt(
            blank_slide, monotone_theme,
            tasks=[{"name": "進行中", "start": 0, "duration": 2, "progress": 0.6}],
            phases=["A", "B"],
            left=Inches(0.5), top=Inches(3),
        )
        count_with_progress = len(blank_slide.shapes) - before2
        assert count_with_progress > count_no_progress


class TestTier3Components:
    def test_pillars_renders(self, blank_slide, monotone_theme):
        before = len(blank_slide.shapes)
        add_pillars(
            blank_slide, monotone_theme,
            items=[
                {"title": "基盤整備", "body": "データ統合\nシステム刷新", "kpi": "3億円"},
                {"title": "業務変革", "body": "自動化\nプロセス改善", "kpi": "1.5億円"},
                {"title": "価値創造", "body": "AI活用\n新規事業", "kpi": "5億円"},
            ],
            left=Inches(0.6), top=Inches(1.5),
        )
        assert len(blank_slide.shapes) > before + 2

    def test_swot_renders(self, blank_slide, monotone_theme):
        before = len(blank_slide.shapes)
        add_swot(
            blank_slide, monotone_theme,
            cells=[
                {"title": "強み", "items": ["ブランド力", "技術力"]},
                {"title": "弱み", "items": ["コスト高", "人材不足"]},
                {"title": "機会", "items": ["市場拡大", "規制緩和"]},
                {"title": "脅威", "items": ["競合参入", "景気後退"]},
            ],
            left=Inches(0.6), top=Inches(1.5),
        )
        assert len(blank_slide.shapes) > before + 3

    def test_heatmap_renders(self, blank_slide, monotone_theme):
        before = len(blank_slide.shapes)
        add_heatmap(
            blank_slide, monotone_theme,
            col_headers=["Q1", "Q2", "Q3", "Q4"],
            row_headers=["北米", "欧州", "アジア"],
            values=[[80, 85, 90, 95], [60, 65, 70, 75], [40, 55, 70, 85]],
            left=Inches(0.6), top=Inches(1.5),
        )
        assert len(blank_slide.shapes) > before + 3

    def test_heatmap_uniform_values(self, blank_slide, monotone_theme):
        add_heatmap(
            blank_slide, monotone_theme,
            col_headers=["A", "B"],
            row_headers=["X"],
            values=[[50, 50]],
            left=Inches(0.6), top=Inches(1.5),
        )

    def test_benchmark_bar_renders(self, blank_slide, monotone_theme):
        before = len(blank_slide.shapes)
        add_benchmark_bar(
            blank_slide, monotone_theme,
            items=[
                {"label": "自社", "value": 85, "is_self": True},
                {"label": "競合A", "value": 72},
                {"label": "競合B", "value": 68},
                {"label": "業界平均", "value": 75},
            ],
            left=Inches(0.6), top=Inches(1.5),
            unit="%",
        )
        assert len(blank_slide.shapes) > before + 3

    def test_benchmark_bar_is_self_highlighted(self, blank_slide, monotone_theme):
        add_benchmark_bar(
            blank_slide, monotone_theme,
            items=[
                {"label": "自社", "value": 100, "is_self": True},
                {"label": "他社", "value": 80},
            ],
            left=Inches(0.6), top=Inches(1.5),
        )
