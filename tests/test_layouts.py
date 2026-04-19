import pytest
from pptx import Presentation
from pptx.util import Inches

from src.themes.monotone import MonotoneTheme
from src.themes.dark import DarkTheme
from src.layouts.cover import CoverLayout
from src.layouts.section_divider import SectionDividerLayout
from src.layouts.agenda import AgendaLayout
from src.layouts.content import ContentLayout
from src.layouts.chart_page import ChartPageLayout
from src.layouts.comparison import ComparisonLayout
from src.layouts.closing import ClosingLayout


@pytest.fixture
def theme():
    return MonotoneTheme()


@pytest.fixture
def dark_theme():
    return DarkTheme()


@pytest.fixture
def prs():
    p = Presentation()
    p.slide_width = Inches(13.333)
    p.slide_height = Inches(7.5)
    return p


def make_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


class TestCoverLayout:
    def test_render(self, prs, theme):
        slide = make_slide(prs)
        layout = CoverLayout()
        layout.render(slide, theme, {
            "title": "DX推進戦略提案書",
            "subtitle": "2026年度計画",
            "client": "株式会社ABC",
            "date": "2026年4月",
        })
        assert len(slide.shapes) >= 2

    def test_render_dark_theme(self, prs, dark_theme):
        slide = make_slide(prs)
        layout = CoverLayout()
        layout.render(slide, dark_theme, {"title": "テスト"})
        assert len(slide.shapes) >= 1


class TestSectionDividerLayout:
    def test_render(self, prs, theme):
        slide = make_slide(prs)
        layout = SectionDividerLayout()
        layout.render(slide, theme, {
            "section_number": 1,
            "section_title": "現状分析",
        })
        assert len(slide.shapes) >= 1


class TestAgendaLayout:
    def test_render(self, prs, theme):
        slide = make_slide(prs)
        layout = AgendaLayout()
        layout.render(slide, theme, {
            "items": ["現状分析", "課題整理", "戦略提案", "実行計画"],
        })
        assert len(slide.shapes) >= 1

    def test_render_with_highlight(self, prs, theme):
        slide = make_slide(prs)
        layout = AgendaLayout()
        layout.render(slide, theme, {
            "items": ["現状分析", "課題整理", "戦略提案", "実行計画"],
            "highlight": 1,
        })
        assert len(slide.shapes) >= 1


class TestContentLayout:
    def test_render_single_column(self, prs, theme):
        slide = make_slide(prs)
        layout = ContentLayout()
        layout.render(slide, theme, {
            "title": "テストタイトル",
            "columns": 1,
            "components": [
                {"type": "bullets", "items": ["項目1", "項目2"]},
            ],
        })
        assert len(slide.shapes) >= 2

    def test_render_two_columns(self, prs, theme):
        slide = make_slide(prs)
        layout = ContentLayout()
        layout.render(slide, theme, {
            "title": "比較",
            "columns": 2,
            "components": [
                {"type": "bullets", "items": ["左1", "左2"]},
                {"type": "bullets", "items": ["右1", "右2"]},
            ],
        })
        assert len(slide.shapes) >= 3


class TestChartPageLayout:
    def test_render_with_key_points(self, prs, theme):
        slide = make_slide(prs)
        layout = ChartPageLayout()
        layout.render(slide, theme, {
            "title": "売上推移",
            "chart": {
                "type": "bar",
                "data": {
                    "labels": ["2023", "2024", "2025"],
                    "series": [{"name": "売上", "values": [100, 150, 220]}],
                },
            },
            "key_points": ["成長率15%", "目標達成"],
        })
        assert len(slide.shapes) >= 3

    def test_render_full_width(self, prs, theme):
        slide = make_slide(prs)
        layout = ChartPageLayout()
        layout.render(slide, theme, {
            "title": "市場分析",
            "chart": {
                "type": "line",
                "data": {
                    "labels": ["Q1", "Q2", "Q3"],
                    "series": [{"name": "値", "values": [10, 20, 30]}],
                },
            },
        })
        assert len(slide.shapes) >= 2


class TestComparisonLayout:
    def test_render(self, prs, theme):
        slide = make_slide(prs)
        layout = ComparisonLayout()
        layout.render(slide, theme, {
            "title": "Before / After",
            "left_title": "Before",
            "left_components": [
                {"type": "bullets", "items": ["旧手法1", "旧手法2"]},
            ],
            "right_title": "After",
            "right_components": [
                {"type": "bullets", "items": ["新手法1", "新手法2"]},
            ],
        })
        assert len(slide.shapes) >= 4


class TestClosingLayout:
    def test_render_summary(self, prs, theme):
        slide = make_slide(prs)
        layout = ClosingLayout()
        layout.render(slide, theme, {
            "summary": ["要点1", "要点2", "要点3"],
            "next_steps": ["ステップ1", "ステップ2"],
        })
        assert len(slide.shapes) >= 2

    def test_render_thank_you(self, prs, theme):
        slide = make_slide(prs)
        layout = ClosingLayout()
        layout.render(slide, theme, {
            "type": "thank_you",
            "contact": "example@company.com",
        })
        assert len(slide.shapes) >= 1
