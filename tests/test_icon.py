import pytest
from pptx import Presentation
from pptx.util import Inches

from src.components.icon import add_icon_with_label, add_kpi_card, add_icon_row
from src.themes.monotone import MonotoneTheme


@pytest.fixture
def theme():
    return MonotoneTheme()


@pytest.fixture
def slide():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)


class TestIconWithLabel:
    def test_adds_shapes(self, slide, theme):
        add_icon_with_label(slide, theme, "check", "完了", Inches(1.0), Inches(1.0))
        assert len(slide.shapes) >= 2


class TestKpiCard:
    def test_adds_shapes(self, slide, theme):
        add_kpi_card(slide, theme, "125", "億円", "年間売上", Inches(1.0), Inches(1.0))
        assert len(slide.shapes) >= 1

    def test_text_content(self, slide, theme):
        add_kpi_card(slide, theme, "125", "億円", "年間売上", Inches(1.0), Inches(1.0))
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for p in shape.text_frame.paragraphs:
                    if p.text.strip():
                        texts.append(p.text.strip())
        assert any("125" in t for t in texts)


class TestIconRow:
    def test_adds_shapes(self, slide, theme):
        items = [
            {"icon": "circle", "label": "項目1"},
            {"icon": "square", "label": "項目2"},
            {"icon": "triangle", "label": "項目3"},
        ]
        add_icon_row(slide, theme, items, Inches(0.6), Inches(2.0))
        assert len(slide.shapes) >= 6
