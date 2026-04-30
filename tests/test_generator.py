import os

import pytest
from pptx import Presentation

from src.generator import generate_pptx


@pytest.fixture
def minimal_config():
    return {
        "theme": "monotone",
        "slides": [
            {
                "layout": "cover",
                "data": {"title": "テスト資料", "subtitle": "サブタイトル"},
            },
        ],
    }


def test_generate_minimal(tmp_path, minimal_config):
    output = tmp_path / "out.pptx"
    prs = generate_pptx(minimal_config, str(output))
    assert output.exists()
    assert len(prs.slides) == 1


def test_generate_all_layouts(tmp_path):
    config = {
        "theme": "monotone",
        "slides": [
            {"layout": "cover", "data": {"title": "表紙"}},
            {"layout": "section_divider", "data": {"section_number": 1, "section_title": "現状"}},
            {"layout": "agenda", "data": {"items": ["A", "B", "C"]}},
            {
                "layout": "content",
                "data": {
                    "title": "内容",
                    "columns": 1,
                    "components": [{"type": "bullets", "items": ["1", "2"]}],
                },
            },
            {
                "layout": "chart_page",
                "data": {
                    "title": "チャート",
                    "chart": {
                        "type": "bar",
                        "data": {
                            "labels": ["Q1", "Q2"],
                            "series": [{"name": "売上", "values": [10, 20]}],
                        },
                    },
                },
            },
            {
                "layout": "comparison",
                "data": {
                    "title": "比較",
                    "left_title": "Before",
                    "left_components": [{"type": "bullets", "items": ["旧"]}],
                    "right_title": "After",
                    "right_components": [{"type": "bullets", "items": ["新"]}],
                },
            },
            {"layout": "closing", "data": {"summary": ["要点"], "next_steps": ["次"]}},
        ],
    }
    output = tmp_path / "all.pptx"
    prs = generate_pptx(config, str(output))
    assert output.exists()
    assert len(prs.slides) == 7


def test_generate_dark_theme(tmp_path, minimal_config):
    minimal_config["theme"] = "dark"
    output = tmp_path / "dark.pptx"
    prs = generate_pptx(minimal_config, str(output))
    assert output.exists()
    assert len(prs.slides) == 1


def test_generate_colorful_theme(tmp_path, minimal_config):
    minimal_config["theme"] = "colorful"
    output = tmp_path / "colorful.pptx"
    prs = generate_pptx(minimal_config, str(output))
    assert output.exists()


def test_generate_unknown_layout_raises(tmp_path):
    config = {
        "theme": "monotone",
        "slides": [{"layout": "nonexistent", "data": {}}],
    }
    # validator が先に未知レイアウトを検出する
    with pytest.raises(ValueError, match="nonexistent"):
        generate_pptx(config, str(tmp_path / "err.pptx"))


def test_generate_unknown_layout_without_validate_raises_from_dispatch(tmp_path):
    config = {
        "theme": "monotone",
        "slides": [{"layout": "nonexistent", "data": {}}],
    }
    with pytest.raises(ValueError, match="Unknown layout"):
        generate_pptx(config, str(tmp_path / "err.pptx"), validate=False)


def test_generate_empty_slides(tmp_path):
    # validator は minItems:1 を要求するため validate=False で実行
    config = {"theme": "monotone", "slides": []}
    output = tmp_path / "empty.pptx"
    prs = generate_pptx(config, str(output), validate=False)
    assert output.exists()
    assert len(prs.slides) == 0


def test_generate_default_theme(tmp_path):
    # theme は必須化されたため、validate=False で旧挙動を確認
    config = {"slides": [{"layout": "cover", "data": {"title": "T"}}]}
    output = tmp_path / "default.pptx"
    prs = generate_pptx(config, str(output), validate=False)
    assert output.exists()
