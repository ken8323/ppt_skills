import pytest
from pptx import Presentation
from pptx.util import Inches

from src.components.shape import (
    add_matrix_2x2, add_pyramid, add_process_flow, add_cycle, add_org_chart
)
from src.themes.monotone import MonotoneTheme


@pytest.fixture
def theme():
    return MonotoneTheme()


@pytest.fixture
def slide():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)


class TestMatrix2x2:
    def test_adds_shapes(self, slide, theme):
        add_matrix_2x2(
            slide, theme,
            x_axis="コスト", y_axis="効果",
            quadrants=["高効果/低コスト", "高効果/高コスト", "低効果/低コスト", "低効果/高コスト"],
            left=Inches(1.0), top=Inches(1.5),
        )
        assert len(slide.shapes) >= 4


class TestPyramid:
    def test_adds_shapes(self, slide, theme):
        add_pyramid(
            slide, theme,
            levels=["戦略", "戦術", "実行"],
            left=Inches(3.0), top=Inches(1.5),
        )
        assert len(slide.shapes) >= 3


class TestProcessFlow:
    def test_adds_shapes(self, slide, theme):
        add_process_flow(
            slide, theme,
            steps=["計画", "設計", "実装", "テスト"],
            left=Inches(0.6), top=Inches(2.0),
        )
        assert len(slide.shapes) >= 4


class TestCycle:
    def test_adds_shapes(self, slide, theme):
        add_cycle(
            slide, theme,
            items=["Plan", "Do", "Check", "Act"],
            left=Inches(3.0), top=Inches(1.5),
        )
        assert len(slide.shapes) >= 4


class TestOrgChart:
    def test_adds_shapes(self, slide, theme):
        add_org_chart(
            slide, theme,
            data={"name": "CEO", "children": [
                {"name": "CTO"},
                {"name": "CFO"},
            ]},
            left=Inches(2.0), top=Inches(1.5),
        )
        assert len(slide.shapes) >= 3
