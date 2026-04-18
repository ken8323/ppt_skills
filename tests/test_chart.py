import pytest
from pptx import Presentation
from pptx.util import Inches

from src.components.chart import add_bar_chart, add_line_chart, add_pie_chart, add_waterfall
from src.themes.monotone import MonotoneTheme


@pytest.fixture
def theme():
    return MonotoneTheme()


@pytest.fixture
def slide():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)


class TestBarChart:
    def test_adds_chart(self, slide, theme):
        data = {
            "labels": ["2023", "2024", "2025"],
            "series": [{"name": "売上", "values": [100, 150, 220]}],
        }
        add_bar_chart(slide, theme, data, Inches(0.6), Inches(1.5))
        assert len(slide.shapes) >= 1

    def test_multi_series(self, slide, theme):
        data = {
            "labels": ["Q1", "Q2", "Q3"],
            "series": [
                {"name": "売上", "values": [100, 150, 220]},
                {"name": "利益", "values": [20, 35, 50]},
            ],
        }
        add_bar_chart(slide, theme, data, Inches(0.6), Inches(1.5))
        assert len(slide.shapes) >= 1

    def test_horizontal_bar(self, slide, theme):
        data = {
            "labels": ["A", "B", "C"],
            "series": [{"name": "値", "values": [10, 20, 30]}],
        }
        add_bar_chart(slide, theme, data, Inches(0.6), Inches(1.5), horizontal=True)
        assert len(slide.shapes) >= 1


class TestLineChart:
    def test_adds_chart(self, slide, theme):
        data = {
            "labels": ["1月", "2月", "3月"],
            "series": [{"name": "推移", "values": [10, 20, 15]}],
        }
        add_line_chart(slide, theme, data, Inches(0.6), Inches(1.5))
        assert len(slide.shapes) >= 1


class TestPieChart:
    def test_adds_chart(self, slide, theme):
        data = {
            "labels": ["製品A", "製品B", "製品C"],
            "values": [40, 35, 25],
        }
        add_pie_chart(slide, theme, data, Inches(0.6), Inches(1.5))
        assert len(slide.shapes) >= 1


class TestWaterfall:
    def test_adds_chart(self, slide, theme):
        data = {
            "labels": ["開始", "+営業", "+開発", "-コスト", "合計"],
            "values": [100, 50, 30, -20, 160],
        }
        add_waterfall(slide, theme, data, Inches(0.6), Inches(1.5))
        assert len(slide.shapes) >= 1
