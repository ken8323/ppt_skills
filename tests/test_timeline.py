import pytest
from pptx import Presentation
from pptx.util import Inches

from src.components.timeline import add_timeline, add_gantt
from src.themes.monotone import MonotoneTheme


@pytest.fixture
def theme():
    return MonotoneTheme()


@pytest.fixture
def slide():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)


class TestTimeline:
    def test_adds_shapes(self, slide, theme):
        milestones = [
            {"date": "2026/4", "label": "キックオフ"},
            {"date": "2026/6", "label": "要件定義完了"},
            {"date": "2026/9", "label": "開発完了"},
            {"date": "2026/12", "label": "リリース"},
        ]
        add_timeline(slide, theme, milestones, Inches(0.6), Inches(2.0))
        assert len(slide.shapes) >= 5


class TestGantt:
    def test_adds_shapes(self, slide, theme):
        tasks = [
            {"name": "要件定義", "start": 0, "duration": 2},
            {"name": "設計", "start": 1, "duration": 3},
            {"name": "実装", "start": 3, "duration": 4},
        ]
        phases = ["4月", "5月", "6月", "7月", "8月", "9月"]
        add_gantt(slide, theme, tasks, phases, Inches(0.6), Inches(1.5))
        assert len(slide.shapes) >= 3
